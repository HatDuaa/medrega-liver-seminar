# -*- coding: utf-8 -*-
"""5 slide bổ sung phần paper Mục 4-6 (VN để review, EN sau). Chèn trước S12 của deck gộp. cwd=repo root."""
import sys, re
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x1A,0x3A,0x5C); SUBTXT=RGBColor(0xB8,0xCF,0xE0); FOOTBG=RGBColor(0xEE,0xF4,0xFB)
FOOTTX=RGBColor(0x6B,0x7F,0x95); CARD2=RGBColor(0xE8,0xF0,0xF8); GREY=RGBColor(0x75,0x70,0x70)
INK=RGBColor(0x1C,0x27,0x36); INFO=RGBColor(0x1F,0x55,0x8C); CAV=RGBColor(0xA0,0x46,0x46)
CAVBG=RGBColor(0xFB,0xF1,0xF1); PAGECL=RGBColor(0x5A,0x5A,0x5A); WHITE=RGBColor(0xFF,0xFF,0xFF)
GREENC=RGBColor(0x1B,0x7A,0x4B); ORANGEC=RGBColor(0xC0,0x6A,0x1B)
TFONT,BFONT="Cambria","Calibri"; SW,SH=10.0,5.625

def build(LANG):
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    BLANK=prs.slide_layouts[6]
    EN=(LANG=="en")
    def T(vn,en): return en if EN else vn
    def _set(run,size,color,bold=False,italic=False,font=BFONT):
        run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic; run.font.name=font; run.font.color.rgb=color
    def rect(slide,l,t,w,h,fill):
        sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False; return sp
    def _runs(p,text):
        for seg in re.split(r"(\*\*.*?\*\*)",text):
            if not seg: continue
            r=p.add_run()
            if seg.startswith("**") and seg.endswith("**"): r.text=seg[2:-2]; yield (r,True)
            else: r.text=seg; yield (r,False)
    def tb_(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
        tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
        tf.word_wrap=True; tf.vertical_anchor=anchor
        tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1); return tf
    def header(slide,title,subtitle=None):
        h=1.30 if subtitle else 0.97; rect(slide,0,0,SW,h,NAVY)
        tf=tb_(slide,0.35,0.12 if subtitle else 0.20,9.3,0.60)
        for r,_ in _runs(tf.paragraphs[0],title): _set(r,22,WHITE,bold=True,font=TFONT)
        if subtitle:
            tf2=tb_(slide,0.37,0.76,9.2,0.45); r=tf2.paragraphs[0].add_run(); r.text=subtitle; _set(r,11.5,SUBTXT,italic=True)
        return h+0.15
    def footer(slide,page):
        rect(slide,0,5.34,SW,0.29,FOOTBG); tf=tb_(slide,0.35,5.40,7,0.18)
        r=tf.paragraphs[0].add_run(); r.text="MedRegA · Wang et al., ICLR 2025"; _set(r,8.5,FOOTTX)
        tf2=tb_(slide,9.0,5.40,0.8,0.18); tf2.paragraphs[0].alignment=PP_ALIGN.RIGHT
        r2=tf2.paragraphs[0].add_run(); r2.text=str(page); _set(r2,9,PAGECL)
    def bullets(slide,items,l,t,w,h,base=14,gap=9):
        tf=tb_(slide,l,t,w,h); first=True
        for text,lvl in items:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.level=lvl; p.space_after=Pt(gap); p.space_before=Pt(0)
            pr=p.add_run(); pr.text=("▪  " if lvl==0 else "–  "); _set(pr,base,NAVY if lvl==0 else GREY,bold=(lvl==0))
            for r,isb in _runs(p,text): _set(r,base if lvl==0 else base-1,INK if lvl==0 else GREY,bold=isb)
    def notebox(slide,label,text,l,t,w,h,kind="info"):
        rect(slide,l,t,w,h,CARD2 if kind=="info" else CAVBG); tf=tb_(slide,l+0.12,t+0.07,w-0.24,h-0.14)
        r=tf.paragraphs[0].add_run(); r.text=label; _set(r,10.5,INFO if kind=="info" else CAV,bold=True)
        p2=tf.add_paragraph(); p2.space_before=Pt(1)
        for r,isb in _runs(p2,text): _set(r,10,INK,bold=isb)
    def image(slide,path,l,t,w=None,h=None,caption=None,cap_fs=8.5):
        kw={}
        if w: kw["width"]=Inches(w)
        if h: kw["height"]=Inches(h)
        pic=slide.shapes.add_picture(path,Inches(l),Inches(t),**kw)
        if caption:
            tf=tb_(slide,l,t+pic.height/914400+0.03,max(pic.width/914400,2.0),0.32)
            pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER; r=pp.add_run(); r.text=caption; _set(r,cap_fs,FOOTTX,italic=True)
        return pic
    def table(slide,headers,rows,l,t,w,colw=None,fs=11,hfs=11,rh=0.34,first_bold=True,hl=None):
        g=slide.shapes.add_table(len(rows)+1,len(headers),Inches(l),Inches(t),Inches(w),Inches(rh*(len(rows)+1))).table
        g.first_row=False; g.horz_banding=False
        if colw:
            for i,cw in enumerate(colw): g.columns[i].width=Inches(cw)
        for j,htx in enumerate(headers):
            c=g.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Pt(5); c.margin_top=Pt(2); c.margin_bottom=Pt(2)
            pp=c.text_frame.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
            for r,isb in _runs(pp,htx): _set(r,hfs,WHITE,bold=True)
        for i,row in enumerate(rows,1):
            for j,val in enumerate(row):
                c=g.cell(i,j); c.fill.solid()
                c.fill.fore_color.rgb=(CARD2 if (hl and i-1 in hl) else (WHITE if i%2 else CARD2))
                c.vertical_anchor=MSO_ANCHOR.MIDDLE
                c.margin_left=Pt(5); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
                pp=c.text_frame.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
                for r,isb in _runs(pp,val): _set(r,fs,INK,bold=(isb or (first_bold and j==0)))
    def notes(slide,text): slide.notes_slide.notes_text_frame.text=text
    def stagebox(slide,l,t,w,h,title,lines):
        rect(slide,l,t,w,h,CARD2); tf=tb_(slide,l+0.14,t+0.09,w-0.28,h-0.18)
        p=tf.paragraphs[0]; r=p.add_run(); r.text=title; _set(r,11,INFO,bold=True)
        for key,val in lines:
            pp=tf.add_paragraph(); pp.space_before=Pt(2.5)
            rk=pp.add_run(); rk.text=key+": "; _set(rk,9.6,NAVY,bold=True)
            for rr,isb in _runs(pp,val): _set(rr,9.6,INK,bold=isb)
    def steppanel(slide,l,t,w,h,title,steps,result):
        rect(slide,l,t,w,h,CARD2); tf=tb_(slide,l+0.16,t+0.11,w-0.32,h-0.22)
        p=tf.paragraphs[0]; r=p.add_run(); r.text=title; _set(r,12.5,INFO,bold=True,font=TFONT)
        for i,stx in enumerate(steps,1):
            pp=tf.add_paragraph(); pp.space_before=Pt(3.5)
            rn=pp.add_run(); rn.text=f"{i}.  "; _set(rn,10.5,NAVY,bold=True)
            for rr,isb in _runs(pp,stx): _set(rr,10.5,INK,bold=isb)
        pr=tf.add_paragraph(); pr.space_before=Pt(5)
        rc=pr.add_run(); rc.text="✓  "+result; _set(rc,10,GREENC,bold=True)

    TRN =T("seminar/code/paper_train_stages_vn.png","seminar/code/paper_train_stages_en.png")
    BAR =T("seminar/code/paper_general_bar_vn.png","seminar/code/paper_general_bar_en.png")
    FIG4="seminar/code/paper_fig4_orig.png"   # Figure 4 gốc từ paper (chung 2 ngôn ngữ)
    FIG5="seminar/code/paper_fig5_orig.png"   # Figure 5 gốc từ paper

    def taskcard(slide,l,t,w,h,name,desc):
        rect(slide,l,t,w,h,CARD2)
        tf=tb_(slide,l+0.12,t+0.10,w-0.24,0.4); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=name; _set(r,13,NAVY,bold=True,font=TFONT)
        tf2=tb_(slide,l+0.12,t+0.52,w-0.24,h-0.6); p2=tf2.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
        for rr,isb in _runs(p2,desc): _set(rr,9.8,INK,bold=isb)

    # ===== P1 — Box-as-text + 3 nhiệm vụ (Figure 4 gốc) =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("Ý tưởng cốt lõi: mã hoá vùng thành chữ (box-as-text)","Region-Centric Task Formulation"),
        T("Toạ độ vùng viết thẳng vào câu → nền cho 3 nhiệm vụ region-centric","Regions written as bounding-box tokens — the basis of the three region-centric tasks"))
    image(s,FIG4,0.80,bt-0.02,w=8.40,caption=T("Figure 4 — ví dụ 3 nhiệm vụ region-centric (paper)","Figure 4 — the three region-centric tasks (paper)"),cap_fs=8)
    notebox(s,T("Mã hoá box-as-text","Box-as-text"),
      T("Toạ độ hộp viết thẳng vào câu: **<ref>tên đối tượng</ref> <box>[x1,y1,x2,y2]</box>**, chuẩn hoá về số nguyên **[0, 1000)** (tức **0–999**) → không cần “đầu ra” riêng cho box; cả 3 nhiệm vụ đều quy về sinh văn bản.",
        "Box coordinates written straight into the sentence: **<ref>object</ref> <box>[x1,y1,x2,y2]</box>**, normalized to integers in **[0, 1000)** (i.e. **0–999**) → no separate box output; all three tasks reduce to text generation."),
      0.40,bt+3.30,9.35,0.64,kind="info")
    footer(s,1)

    # ===== P2 — Huấn luyện 2 giai đoạn (chi tiết mục tiêu + dữ liệu) =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("Huấn luyện 2 giai đoạn","Model Training — two steps"),
        T("Nền InternVL 1.2 (InternViT-6B + Yi-34B) — căn chỉnh trước, rồi tinh chỉnh theo chỉ dẫn",
          "Base InternVL 1.2 (InternViT-6B + Yi-34B) — alignment training, then instruction tuning"))
    image(s,TRN,1.75,bt-0.02,w=6.5)
    sy=bt+1.50; sh=1.36
    stagebox(s,0.40,sy,4.55,sh,T("Giai đoạn 1 — Alignment","Stage 1 — Alignment training"),[
      (T("Mục tiêu","Goal"),T("dạy **cầu nối** phiên dịch ảnh y khoa → ngôn ngữ","teach the **connector** to map medical images → language")),
      (T("Dữ liệu","Data"),T("ảnh y khoa + **chú thích** (captioning)","medical images + **captions** (captioning)")),
      (T("Huấn luyện","Train"),T("chỉ **cầu nối**; vision & LLM đóng băng","**connector** only; vision & LLM frozen")),
    ])
    stagebox(s,5.05,sy,4.70,sh,T("Giai đoạn 2 — Instruction tuning","Stage 2 — Instruction tuning"),[
      (T("Mục tiêu","Goal"),T("làm nhiều tác vụ theo **chỉ dẫn**","perform many tasks from **instructions**")),
      (T("Dữ liệu","Data"),T("public (VQA, báo cáo, phân loại) + **MedRegInstruct** (3 tác vụ vùng)","public (VQA, report, classification) + **MedRegInstruct** (3 region tasks)")),
      (T("Huấn luyện","Train"),T("**LLM**; loss = dự đoán token kế tiếp","**LLM**; loss = next-token prediction")),
    ])
    notebox(s,T("MedRegInstruct — dữ liệu vùng song ngữ (giai đoạn 2 học từ đây)","MedRegInstruct — bilingual region data (used in Stage 2)"),
      T("**790K** mẫu tự dựng bán tự động: **550K Region-Text** (SA-Med2D-20M) + **240K Region-Grounded** (MIMIC-CXR + bệnh viện TQ). Báo cáo **tiếng Trung** từ bệnh viện → năng lực **song ngữ Anh–Trung** mà baseline khác thiếu.",
        "**790K** semi-automatically built samples: **550K Region-Text** (SA-Med2D-20M) + **240K Region-Grounded** (MIMIC-CXR + a Chinese hospital). **Chinese** hospital reports → the **English–Chinese** ability other baselines lack."),
      0.40,sy+sh+0.10,9.35,0.74,kind="info")
    footer(s,2)

    # ===== P2b — GIAI ĐOẠN 2: 3 KIỂU CHỈ DẪN (nhiệm vụ vùng) =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("Giai đoạn 2 dạy gì? — 3 kiểu “chỉ dẫn” vùng","Stage-2 instructions — the three region-centric tasks"),
        T("Mỗi nhiệm vụ là một kiểu chỉ dẫn: đưa vào gì → model trả gì",
          "Each task is an instruction type: what goes in → what the model returns"))
    table(s,[T("Nhiệm vụ","Task"),T("Đưa vào (input)","Input"),T("Model trả (output)","Model returns")],
      [[T("**Region → Text**","**Region-to-Text Identification**"),
        T("ảnh + **box [toạ độ]**","image + **box [coords]**"),
        T("**tên vùng** — vd “gan”, “khối u”","**region name** — e.g. “liver”, “tumor”")],
       [T("**Text → Region**","**Text-to-Region Detection**"),
        T("ảnh + **tên** — vd “gan ở đâu?”","image + **name** — e.g. “where is the liver?”"),
        T("**box [toạ độ]** của vùng","**box [coords]** of the region")],
       [T("**Grounded Report**","**Grounded Report Generation**"),
        T("ảnh (cả ảnh)","the whole image"),
        T("**báo cáo**, mỗi câu gắn 1 box","a **report**, each line tied to a box")]],
      0.40,bt+0.12,9.35,colw=[2.35,3.35,3.65],fs=12,hfs=12,rh=0.62,first_bold=True)
    notebox(s,T("Cách viết “chỉ dẫn”","How instructions are written"),
      T("Toạ độ box viết thẳng trong câu: **<box>[x1,y1,x2,y2]</box>**, chuẩn hoá **0–999** → mọi câu trả lời (kể cả vị trí) đều là **sinh văn bản**. Các cặp (chỉ dẫn → đáp án) lấy từ bộ **MedRegInstruct** (slide trước).",
        "Coordinates are written inline: **<box>[x1,y1,x2,y2]</box>**, normalized to **0–999** → every answer (including location) is **text generation**. These (instruction → answer) pairs come from **MedRegInstruct** (previous slide)."),
      0.40,bt+2.72,9.35,0.80,kind="info")
    footer(s,3)

    # ===== P3 — Kết quả general =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("Kết quả — tác vụ tổng quát","Performance on General Medical Tasks"),
        T("Ngang/nhỉnh baseline ở tiếng Anh, vượt trội ở tiếng Trung","On par in English, clearly ahead in Chinese"))
    image(s,BAR,0.30,bt+0.05,w=5.55)
    bullets(s,[
      (T("**Tiếng Anh:** nhỉnh nhẹ — VQA 67.7 vs MedDr 65.7.","**English:** slight edge — VQA 67.7 vs MedDr 65.7."),0),
      (T("**Tiếng Trung: bứt phá** — Report +28.8, VQA +23 BLEU-1.","**Chinese: big jump** — Report +28.8, VQA +23 BLEU-1."),0),
      (T("**Phân loại ảnh:** F1 47.97 vs 32.65 (+15).","**Classification:** F1 47.97 vs 32.65 (+15)."),0),
      (T("**Regional CoT** (phát hiện vùng trước → suy luận sau): +31.6 F1 (VinDr-SpineXR) → khoanh vùng giúp cả bài toán chung.",
         "**Regional CoT** (detect regions first → then reason): +31.6 F1 (VinDr-SpineXR) → localizing helps even general tasks."),0),
    ],6.05,bt+0.15,3.7,3.0,base=11.5,gap=10)
    notebox(s,T("Chốt","Takeaway"),
      T("Trên sân chung MedRegA hơn nhẹ; bứt phá ở song ngữ và khi thêm suy luận theo vùng.",
        "On general tasks MedRegA is slightly ahead; it pulls away on bilingual data and with region-guided reasoning."),
      0.30,bt+3.25,9.4,0.62,kind="info")
    footer(s,4)

    # ===== P4 — Kết quả region-centric =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("Kết quả — tác vụ Region-Centric (thế mạnh)","Region-Aligned Evaluation on Region-Centric Tasks"),
        T("Baseline gần như bằng 0 — đây là năng lực riêng của MedRegA","Baselines score near 0 — this is MedRegA's unique ability"))
    # left table: Region -> Text (contrast)
    tf=tb_(s,0.40,bt+0.0,4.6,0.3); r=tf.paragraphs[0].add_run()
    r.text=T("Region → Text (gọi tên vùng)","Region → Text (naming)"); _set(r,12,NAVY,bold=True,font=TFONT)
    table(s,[T("Metric","Metric"),T("Baseline","Baseline"),"MedRegA"],
      [["BLEU-1","≤0.75","**69.72**"],["F1","≤1.27","**70.43**"],["BertScore","50.28","**87.13**"]],
      0.40,bt+0.36,4.55,colw=[1.85,1.35,1.35],fs=11,hfs=11,rh=0.42,first_bold=True)
    # right table: Text -> Region (IoU highlight)
    tf=tb_(s,5.15,bt+0.0,4.6,0.3); r=tf.paragraphs[0].add_run()
    r.text=T("Text → Region (vẽ hộp)","Text → Region (detection)"); _set(r,12,NAVY,bold=True,font=TFONT)
    table(s,[T("Metric","Metric"),"InternVL","MedRegA"],
      [["Object F1","56.60","**77.93**"],["Region F1","6.70","**38.24**"],
       ["Alignment F1","5.45","**36.53**"],["**IoU**","12.28","**23.43**"]],
      5.15,bt+0.36,4.55,colw=[1.85,1.35,1.35],fs=11,hfs=11,rh=0.42,first_bold=True,hl=[3])
    # metric ladder
    notebox(s,T("Thang metric (dễ tăng dần)","Metric ladder (harder →)"),
      T("Object = đúng **tên** → Region = đúng **chỗ** → Alignment = đúng **cả hai** → IoU = vẽ **khít** cỡ nào.",
        "Object = right **name** → Region = right **place** → Alignment = **both** → IoU = how **tight**."),
      0.40,bt+2.42,9.35,0.62,kind="info")
    notebox(s,T("Lưu ý IoU","IoU note"),
      T("IoU 23.43 tính bằng Hungarian matching, chỉ trên các cặp hộp đã ghép. Grounded report IoU 52.07 cao hơn vì đích là **cơ quan lớn** (dễ), còn detection nhắm **khối u nhỏ** (khó).",
        "IoU 23.43 is computed with Hungarian matching over matched pairs only. Grounded-report IoU 52.07 is higher because targets are **large organs** (easier), while detection aims at **small tumors** (harder)."),
      0.40,bt+3.16,9.35,0.72,kind="caveat")
    footer(s,5)

    # ===== P5 — Chốt paper + cầu nối =====
    s=prs.slides.add_slide(BLANK); bt=header(s,
        T("MedRegA — chốt lại & cầu nối","Conclusion — and the bridge to our work"),
        T("Paper đóng góp gì, còn hạn chế gì (nhóm nhận xét)","What the paper contributes, and its limits (our observations)"))
    notebox(s,T("Đóng góp của paper","Paper's contributions"),
      T("Bộ 3 nhiệm vụ Region-Centric + dữ liệu **MedRegInstruct** (gán nhãn bán tự động) + mô hình **MedRegA** region-aware, song ngữ. Giá trị cốt lõi: gắn **vùng ↔ văn bản** → tăng khả năng diễn giải & tương tác lâm sàng.",
        "Three region-centric tasks + the **MedRegInstruct** dataset (semi-automatic labeling) + the **MedRegA** region-aware bilingual model. Core value: linking **region ↔ text** → more interpretable & clinically interactive."),
      0.40,bt+0.05,9.35,0.98,kind="info")
    bullets(s,[
      (T("**Chỉ song ngữ Trung–Anh** — chưa có tiếng Việt.","**Only Chinese–English** — no Vietnamese."),0),
      (T("**Chủ yếu ảnh 2D** — chưa khai thác 3D / CT đa pha.","**Mostly 2D** — no 3D / multi-phase CT."),0),
      (T("**Chưa external validation / calibration** — độ tin cậy triển khai còn để ngỏ.","**No external validation / calibration** — deployment reliability still open."),0),
    ],0.40,bt+1.42,9.35,1.5,base=13,gap=9)
    notebox(s,T("→ Cầu nối","→ Bridge"),
      T("MedRegA mạnh ở vùng-ngôn ngữ & diễn giải, nhưng còn hạn ở tiếng Việt, dữ liệu 3D và kiểm định ngoài → đó là chỗ **nhóm mình bắt tay cải tiến**.",
        "MedRegA is strong on region-language & interpretability, but limited on Vietnamese, 3D data, and external validation → that's where **our work steps in**."),
      0.40,bt+2.95,9.35,0.72,kind="info")
    tf=tb_(s,0.40,bt+3.74,9.35,0.28); r=tf.paragraphs[0].add_run()
    r.text=T("* Paper không có mục Hạn chế/Hướng phát triển riêng — đây là nhận xét thêm của nhóm.",
             "* The paper has no explicit Limitations/Future-Work section — these are our own observations.")
    _set(r,8.5,GREY,italic=True)
    footer(s,6)

    out=f"seminar/cai_tien/slides_paper_sec4to6_{'EN' if EN else 'VN'}.pptx"
    prs.save(out); print("SAVED",out,"| slides =",len(prs.slides._sldIdLst))

for lang in ("vn","en"): build(lang)
