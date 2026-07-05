# -*- coding: utf-8 -*-
"""Condensed 5-minute standalone deck (our part only), VN + EN. cwd = repo root."""
import sys, re
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x1A,0x3A,0x5C); SUBTXT=RGBColor(0xB8,0xCF,0xE0); FOOTBG=RGBColor(0xEE,0xF4,0xFB)
FOOTTX=RGBColor(0x6B,0x7F,0x95); CARD2=RGBColor(0xE8,0xF0,0xF8); GREY=RGBColor(0x75,0x70,0x70)
INK=RGBColor(0x1C,0x27,0x36); INFO=RGBColor(0x1F,0x55,0x8C); CAV=RGBColor(0xA0,0x46,0x46)
CAVBG=RGBColor(0xFB,0xF1,0xF1); PAGECL=RGBColor(0x5A,0x5A,0x5A); WHITE=RGBColor(0xFF,0xFF,0xFF)
TFONT,BFONT="Cambria","Calibri"; SW,SH=10.0,5.625

def build(LANG):
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    BLANK=prs.slide_layouts[6]
    EN=(LANG=="en")
    def T(vn,en): return en if EN else vn
    def _set(run,size,color,bold=False,italic=False,font=BFONT):
        run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic; run.font.name=font; run.font.color.rgb=color
    def rect(slide,l,t,w,h,fill):
        sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
        sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False; return sp
    def _runs(p,text):
        for seg in re.split(r"(\*\*.*?\*\*)",text):
            if not seg: continue
            r=p.add_run()
            if seg.startswith("**") and seg.endswith("**"): r.text=seg[2:-2]; yield (r,True)
            else: r.text=seg; yield (r,False)
    def tb_(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
        tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
        tf.word_wrap=True; tf.vertical_anchor=anchor
        tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1); return tf
    def header(slide,title,subtitle=None):
        h=1.30 if subtitle else 0.97; rect(slide,0,0,SW,h,NAVY)
        tf=tb_(slide,0.35,0.12 if subtitle else 0.20,9.3,0.60)
        for r,_ in _runs(tf.paragraphs[0],title): _set(r,22,WHITE,bold=True,font=TFONT)
        if subtitle:
            tf2=tb_(slide,0.37,0.76,9.2,0.45); r=tf2.paragraphs[0].add_run(); r.text=subtitle; _set(r,11.5,SUBTXT,italic=True)
        return h+0.15
    def footer(slide,page):
        rect(slide,0,5.34,SW,0.29,FOOTBG); tf=tb_(slide,0.35,5.40,7,0.18)
        r=tf.paragraphs[0].add_run(); r.text=T("Cải tiến MedRegA · Gemma 4 E4B (u gan LiTS)","MedRegA Improvements · Gemma 4 E4B (liver tumor, LiTS)"); _set(r,8.5,FOOTTX)
        tf2=tb_(slide,9.0,5.40,0.8,0.18); tf2.paragraphs[0].alignment=PP_ALIGN.RIGHT
        r2=tf2.paragraphs[0].add_run(); r2.text=str(page); _set(r2,9,PAGECL)
    def bullets(slide,items,l,t,w,h,base=14,gap=9):
        tf=tb_(slide,l,t,w,h); first=True
        for text,lvl in items:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.level=lvl; p.space_after=Pt(gap); p.space_before=Pt(0)
            pr=p.add_run(); pr.text=("▪  " if lvl==0 else "–  "); _set(pr,base,NAVY if lvl==0 else GREY,bold=(lvl==0))
            for r,isb in _runs(p,text): _set(r,base if lvl==0 else base-1,INK if lvl==0 else GREY,bold=isb)
    def notebox(slide,label,text,l,t,w,h,kind="info"):
        rect(slide,l,t,w,h,CARD2 if kind=="info" else CAVBG); tf=tb_(slide,l+0.12,t+0.07,w-0.24,h-0.14)
        r=tf.paragraphs[0].add_run(); r.text=label; _set(r,10.5,INFO if kind=="info" else CAV,bold=True)
        p2=tf.add_paragraph(); p2.space_before=Pt(1)
        for r,isb in _runs(p2,text): _set(r,10,INK,bold=isb)
    def image(slide,path,l,t,w=None,h=None,caption=None):
        kw={}
        if w: kw["width"]=Inches(w)
        if h: kw["height"]=Inches(h)
        pic=slide.shapes.add_picture(path,Inches(l),Inches(t),**kw)
        if caption:
            tf=tb_(slide,l,t+pic.height/914400+0.03,max(pic.width/914400,2.0),0.3)
            pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER; r=pp.add_run(); r.text=caption; _set(r,8,FOOTTX,italic=True)
        return pic
    def table(slide,headers,rows,l,t,w,colw=None,fs=11,hfs=11,rh=0.34,first_bold=True):
        g=slide.shapes.add_table(len(rows)+1,len(headers),Inches(l),Inches(t),Inches(w),Inches(rh*(len(rows)+1))).table
        g.first_row=False; g.horz_banding=False
        if colw:
            for i,cw in enumerate(colw): g.columns[i].width=Inches(cw)
        for j,htx in enumerate(headers):
            c=g.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY; c.vertical_anchor=MSO_ANCHOR.MIDDLE
            c.margin_left=Pt(5); c.margin_top=Pt(2); c.margin_bottom=Pt(2)
            for r,isb in _runs(c.text_frame.paragraphs[0],htx): _set(r,hfs,WHITE,bold=True)
        for i,row in enumerate(rows,1):
            for j,val in enumerate(row):
                c=g.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2 else CARD2; c.vertical_anchor=MSO_ANCHOR.MIDDLE
                c.margin_left=Pt(5); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
                for r,isb in _runs(c.text_frame.paragraphs[0],val): _set(r,fs,INK,bold=(isb or (first_bold and j==0)))
    IOU=T("seminar/code/iou_frozen_vs_unfrozen.png","seminar/code/iou_frozen_vs_unfrozen_en.png")
    DET=T("seminar/code/detect_before_after.png","seminar/code/detect_before_after_en.png")

    def statcard(slide,l,t,w,h,num,label):
        rect(slide,l,t,w,h,CARD2)
        tf=tb_(slide,l,t+0.14,w,h*0.52,MSO_ANCHOR.BOTTOM); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=num; _set(r,32,NAVY,bold=True,font=TFONT)
        tf2=tb_(slide,l+0.1,t+h*0.58,w-0.2,h*0.4); p2=tf2.paragraphs[0]; p2.alignment=PP_ALIGN.CENTER
        for rr,isb in _runs(p2,label): _set(rr,11,INK,bold=isb)

    # ---- S1 TITLE ----
    s=prs.slides.add_slide(BLANK)
    rect(s,0,0,SW,0.24,NAVY); rect(s,0,1.55,SW,2.30,NAVY)
    tf=tb_(s,0.5,1.72,9.0,1.0,MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=T("Cải tiến của nhóm trên MedRegA","Our Improvements on MedRegA"); _set(r,36,WHITE,bold=True,font=TFONT)
    tf2=tb_(s,0.5,2.82,9.0,0.9);
    for i,line in enumerate([T("Fine-tune Gemma 4 E4B phát hiện u gan trên CT (LiTS)","Fine-tuning Gemma 4 E4B for liver-tumor detection on CT (LiTS)"),
                             T("Mô hình nhỏ — trung thực về đánh giá & độ tin cậy","A small model — honest about evaluation & reliability")]):
        pp=tf2.paragraphs[0] if i==0 else tf2.add_paragraph(); pp.alignment=PP_ALIGN.CENTER
        rr=pp.add_run(); rr.text=line; _set(rr,15,SUBTXT)
    rect(s,0,5.34,SW,0.29,FOOTBG); tf3=tb_(s,0.35,5.40,9,0.18)
    r3=tf3.paragraphs[0].add_run(); r3.text=T("Seminar Học sâu · 2026","Deep Learning Seminar · 2026"); _set(r3,9,FOOTTX)

    # ---- S2 ARCHITECTURE & METHOD (limitation ↔ our choice) ----
    s=prs.slides.add_slide(BLANK); bt=header(s,T("Kiến trúc & cách làm của nhóm","Our architecture & method"),
        T("Mỗi lựa chọn vá đúng MỘT hạn chế của MedRegA","Each choice fixes ONE MedRegA limitation"))
    table(s,[T("Hạn chế của MedRegA","MedRegA limitation"),T("Lựa chọn của nhóm (vì sao)","Our choice (why)")],
      [[T("Dồn lên **LLM 40B**; “mắt” & “cầu” bất động → định vị không sắc",
          "Load on the **40B LLM**; “eyes” & “bridge” frozen → localization not sharp"),
        T("**Gemma 4 E4B (~8B)** + **mở khoá vision** (pIoU 0.015→0.27)",
          "**Gemma 4 E4B (~8B)** + **unlock vision** (pIoU 0.015→0.27)")],
       [T("Chỉ **1 lát trung tâm**, **không ca âm**, box gộp thô",
          "Only **1 central slice**, **no negatives**, coarse box"),
        T("**Đa lát + ca âm + multi-box** (tách từng ổ); split theo bệnh nhân",
          "**Multi-slice + negatives + multi-box** (per-lesion); split by patient")],
       [T("Chi phí **16×H800** — ngoài tầm sinh viên",
          "Cost **16×H800** — out of reach for a student"),
        T("**LoRA + full-finetune vision**, chạy **1 GPU**",
          "**LoRA + full fine-tune vision**, runs on **1 GPU**")],
       [T("**Per-image** (pseudoreplication); không abstain / CI",
          "**Per-image** (pseudoreplication); no abstain / CI"),
        T("**Per-patient + CI**, tách detect/localize, **selective prediction**",
          "**Per-patient + CI**, detect/localize split, **selective prediction**")]],
      0.4,bt+0.1,9.2,colw=[4.4,4.8],fs=11.5,hfs=12,rh=0.62,first_bold=False)
    footer(s,2)

    # ---- S3 OUR RESULTS (publish metrics) ----
    s=prs.slides.add_slide(BLANK); bt=header(s,T("Kết quả của nhóm","Our results"),
        T("Model của nhóm đạt được — công bố trực tiếp","What our model achieves — reported directly"))
    cy=bt+0.12; cw=2.95; gp=0.18
    statcard(s,0.4,cy,cw,1.35,"0.89",T("Detection F1 (bệnh nhân)","Detection F1 (patient)"))
    statcard(s,0.4+cw+gp,cy,cw,1.35,"80%",T("Sensitivity — bắt được u","Sensitivity — catch tumors"))
    statcard(s,0.4+2*(cw+gp),cy,cw,1.35,"0.27",T("Localization pIoU (định vị)","Localization pIoU"))
    bullets(s,[
      (T("**recall@IoU.25 = 54%** · **False-positive = 1.5%** (hiếm khi bịa box) · pIoU CI95 [0.19, 0.36].",
         "**recall@IoU.25 = 54%** · **False-positive = 1.5%** (rarely invents boxes) · pIoU CI95 [0.19, 0.36]."),0),
      (T("**Không quên:** vẫn hội thoại / JSON / thơ / toán (LoRA giữ trọng số gốc).",
         "**No forgetting:** still chats / JSON / poem / math (LoRA keeps base weights)."),0),
      (T("Định vị ~0.27 ≈ **cùng khoảng MedRegA (~0.23)** dù model nhỏ 5× — khác dataset, chỉ tham khảo.",
         "Localization ~0.27 ≈ **same ballpark as MedRegA (~0.23)** despite 5× smaller — different dataset, reference only."),0),
    ],0.4,cy+1.52,9.2,1.5,base=12.5,gap=8)
    notebox(s,T("⚠ Lưu ý","⚠ Note"),T("n nhỏ (~25 bn dương / 2 âm) → CI rộng; đây là PoC, không tuyên bố lâm sàng.",
              "small n (~25 pos patients / 2 neg) → wide CI; a PoC, no clinical claim."),0.4,cy+3.0,9.2,0.6,kind="caveat")
    notes(s,T(
"""GIẢI THÍCH METRIC (Presenter View):
• Detection F1 0.89 = cân bằng recall (bắt u) & precision (không báo bừa); tính PER-PATIENT để không thổi phồng.
• Sensitivity 80% = bắt đúng 20/25 bệnh nhân có u (sót 5) — con số lâm sàng quan trọng nhất (sót u = nguy hiểm).
• Localization pIoU 0.27 = độ chồng lấn box (giao/hợp), CÓ PHẠT khi vẽ thừa/thiếu ổ → mức trung bình, chỗ còn yếu (khoanh được vùng nhưng chưa khít).
• recall@IoU.25 54% = 54% u được khoanh tạm trúng (chồng lấn ≥ 0.25).
• False-positive 1.5% = trên lát KHÔNG u, hiếm khi bịa box (điểm tốt).
• CI95 [0.19, 0.36] = khoảng tin cậy 95%, số dao động vì mẫu nhỏ (không phải số cứng).
• Không quên = nhờ LoRA giữ trọng số gốc → vẫn hội thoại / JSON / thơ / toán.
• ~0.27 vs MedRegA ~0.23 = cùng khoảng dù nhỏ 5×, KHÁC dataset nên chỉ tham khảo, KHÔNG claim hơn.
• Caveat: n nhỏ (25 dương / 2 âm), CI rộng, là PoC — không tuyên bố lâm sàng.

CÂU NÓI (~40s): "Model đạt F1 0.89 ở phát hiện, độ nhạy 80% (20/25 bệnh nhân). Định vị pIoU 0.27 mới trung bình, chúng em không giấu. Điểm cộng: hiếm bịa u (FP 1.5%) và không quên nhờ LoRA. Định vị ngang MedRegA dù nhỏ 5× nhưng khác dữ liệu nên chỉ tham khảo. Đây là PoC, mẫu nhỏ nên còn hạn chế thống kê."
""",
"""METRIC EXPLANATION (Presenter View):
• Detection F1 0.89 = balance of recall (catching tumors) & precision (not over-flagging); computed PER-PATIENT to avoid inflation.
• Sensitivity 80% = caught 20/25 tumor patients (missed 5) — the most clinically important number (missing a tumor is dangerous).
• Localization pIoU 0.27 = box overlap (intersection/union), PENALIZED for extra/missing lesions → moderate, the weaker point (finds the region but not tightly).
• recall@IoU.25 54% = 54% of tumors roughly localized (overlap ≥ 0.25).
• False-positive 1.5% = on tumor-free slices, rarely invents a box (good).
• CI95 [0.19, 0.36] = 95% confidence interval; the number varies because n is small (not a hard value).
• No forgetting = thanks to LoRA keeping base weights → still chats / JSON / poem / math.
• ~0.27 vs MedRegA ~0.23 = same ballpark despite 5× smaller, DIFFERENT dataset so reference only, NOT a 'we beat them' claim.
• Caveat: small n (25 pos / 2 neg), wide CI, a PoC — no clinical claim.

SCRIPT (~40s): "Our model reaches F1 0.89 on detection, 80% sensitivity (20/25 patients). Localization pIoU is 0.27 — only moderate, we don't hide it. Pluses: rarely invents tumors (FP 1.5%) and no forgetting thanks to LoRA. Localization is on par with MedRegA despite being 5× smaller, but different data so reference only. This is a PoC, small sample, so statistical limits remain."
"""))
    footer(s,3)

    # ---- S4 WHY: FINE-TUNE DRIVES IT (baseline ablation) ----
    s=prs.slides.add_slide(BLANK); bt=header(s,T("Vì đâu có kết quả? — nhờ fine-tune","Where do the gains come from? — fine-tuning"),
        T("Không chỉ đổi model: zero-shot (model mới, chưa train) → sau fine-tune","Not just a model swap: zero-shot (new model, untrained) → fine-tuned"))
    bullets(s,[
      (T("**Đổi model thôi KHÔNG đủ:** Gemma zero-shot (chưa fine-tune) chỉ **F1 0.33**, gần như luôn trả “No tumor”.",
         "**Swapping models alone is NOT enough:** zero-shot Gemma (untrained) scores only **F1 0.33**, almost always says “No tumor”."),0),
      (T("**Fine-tune tạo ra năng lực:** F1 0.33→**0.89**, recall 0→**54%**, pIoU 0.002→**0.27**.",
         "**Fine-tuning creates the capability:** F1 0.33→**0.89**, recall 0→**54%**, pIoU 0.002→**0.27**."),0),
      (T("**Bằng chứng thêm:** đóng băng vision → pIoU 0.015; **mở khoá + train → 0.27** — cách train mới là mấu chốt.",
         "**Extra evidence:** frozen vision → pIoU 0.015; **unlock + train → 0.27** — the training choice is the key."),0),
    ],0.4,bt+0.1,5.5,3.0,base=13,gap=12)
    image(s,DET,6.1,bt+0.2,w=3.7,caption=T("Trước → sau fine-tune","Before → after fine-tune"))
    notebox(s,T("Chốt","Takeaway"),T("Kết quả đến từ FINE-TUNE (dữ liệu + mở vision + eval), không chỉ vì đổi sang model mới.",
              "The gains come from FINE-TUNING (data + unlocked vision + eval), not from a newer model."),0.4,bt+3.0,5.5,0.72,kind="info")
    footer(s,4)

    # ---- S5 CONCLUSION ----
    s=prs.slides.add_slide(BLANK); bt=header(s,T("Kết luận","Conclusion"))
    bullets(s,[
      (T("**Chứng minh cơ chế:** region-centric (khoanh-rồi-đọc) khả thi ở **~8B, 1 GPU**.",
         "**Proof of mechanism:** region-centric (localize-then-read) is feasible at **~8B, 1 GPU**."),0),
      (T("**Giá trị:** khung đánh giá trung thực + selective prediction — không phải điểm số.",
         "**Value:** honest evaluation + selective prediction — not the score."),0),
      (T("**Trung thực:** PoC trên 1 dataset / 1 task, n nhỏ, CI rộng.",
         "**Honest:** PoC on 1 dataset / 1 task, small n, wide CI."),0),
      (T("**Tương lai:** ảnh giàu hơn (multi-window, 2.5D) + external validation.",
         "**Future:** richer images (multi-window, 2.5D) + external validation."),0),
    ],0.4,bt+0.1,9.2,3.0,base=15,gap=13)
    notebox(s,T("Câu chốt","Closing"),
      T("Không hứa đã tới đích — chỉ ra con đường và nói rõ chỗ nào còn dốc.",
        "We don't claim to have arrived — we map the road and say where it's still steep."),0.4,bt+3.05,9.2,0.6,kind="info")
    footer(s,5)

    out=f"seminar/cai_tien/slides_5min_{'EN' if EN else 'VN'}.pptx"
    prs.save(out); print("SAVED",out,"| slides =",len(prs.slides._sldIdLst))

for lang in ("vn","en"): build(lang)
