# -*- coding: utf-8 -*-
"""English version: charts EN + deck EN appended to the English paper deck. cwd = repo root."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
import numpy as np
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "seminar/code"
NAVY_H="#1f3a5f"; GREY_H="#9aa7b4"; RED_H="#c0392b"; GREEN_H="#2e7d32"
plt.rcParams["font.family"]="DejaVu Sans"

# ================= ENGLISH CHARTS =================
# detect before/after
fig,ax=plt.subplots(figsize=(7.2,4.2),dpi=150)
labels=["Detection F1\n(patient)","Sensitivity\n(catch tumor)","False-positive\n(neg. slice)"]
zs=[0.33,0.20,0.022]; ft=[0.89,0.80,0.015]; x=np.arange(3); w=0.36
b1=ax.bar(x-w/2,zs,w,label="Zero-shot (untrained)",color=GREY_H)
b2=ax.bar(x+w/2,ft,w,label="After fine-tune (ours)",color=NAVY_H)
for b,v in zip(b1,zs): ax.text(b.get_x()+b.get_width()/2,v+0.015,f"{v:.2f}" if v>=0.1 else f"{v*100:.1f}%",ha="center",fontsize=9,color="#444")
for b,v in zip(b2,ft): ax.text(b.get_x()+b.get_width()/2,v+0.015,f"{v:.2f}" if v>=0.1 else f"{v*100:.1f}%",ha="center",fontsize=9,fontweight="bold",color=NAVY_H)
ax.set_ylim(0,1.0); ax.set_ylabel("Value (0–1)")
ax.set_title("Before → after fine-tune  (n=25 pos / 2 neg)",fontweight="bold",color=NAVY_H)
ax.set_xticks(x); ax.set_xticklabels(labels,fontsize=9)
ax.legend(fontsize=9,loc="upper right",framealpha=0.9)
ax.text(2,0.13,"↓ lower = better",ha="center",fontsize=8,style="italic",color=RED_H)
ax.spines[["top","right"]].set_visible(False)
plt.tight_layout(); plt.savefig(f"{OUT}/detect_before_after_en.png",bbox_inches="tight"); plt.close()

# iou frozen vs unfrozen
fig,ax=plt.subplots(figsize=(5.2,4.2),dpi=150)
cats=["Frozen\nvision","Unlocked\nvision (fine-tune)"]; vals=[0.015,0.27]
bars=ax.bar(cats,vals,color=[GREY_H,NAVY_H],width=0.55)
for b,v in zip(bars,vals): ax.text(b.get_x()+b.get_width()/2,v+0.006,f"{v:.3f}",ha="center",fontsize=11,fontweight="bold",color=NAVY_H)
ax.set_ylim(0,0.32); ax.set_ylabel("pIoU (per-patient, penalized)")
ax.set_title("Unlocking the vision tower = the key step",fontweight="bold",color=NAVY_H)
ax.annotate("",xy=(1,0.27),xytext=(0,0.045),arrowprops=dict(arrowstyle="->",color=GREEN_H,lw=1.8))
ax.text(0.5,0.17,"×18",ha="center",fontsize=12,fontweight="bold",color=GREEN_H)
ax.text(0.5,-0.055,"1 config for illustration — not multi-seed",ha="center",fontsize=7.5,style="italic",color="#666")
ax.spines[["top","right"]].set_visible(False)
plt.tight_layout(); plt.savefig(f"{OUT}/iou_frozen_vs_unfrozen_en.png",bbox_inches="tight"); plt.close()

# bce collapse
fig,ax=plt.subplots(figsize=(4.6,3.5),dpi=150)
lb=["mean-IoU","recall@0.25"]; v3=[0.32,0.54]; v4=[0.071,0.10]; x=np.arange(2); w=0.36
b1=ax.bar(x-w/2,v3,w,label="v3 (good)",color=NAVY_H); b2=ax.bar(x+w/2,v4,w,label="v4 BCE-head (collapse)",color=RED_H)
for b,v in zip(b1,v3): ax.text(b.get_x()+b.get_width()/2,v+0.014,f"{v:.2f}",ha="center",fontsize=9,fontweight="bold",color=NAVY_H)
for b,v in zip(b2,v4): ax.text(b.get_x()+b.get_width()/2,v+0.014,f"{v:.3f}",ha="center",fontsize=9,fontweight="bold",color=RED_H)
ax.set_ylim(0,0.72); ax.set_ylabel("value"); ax.set_xticks(x); ax.set_xticklabels(lb)
ax.set_title("v4 (added BCE-head) COLLAPSES vs v3",fontweight="bold",color=NAVY_H,fontsize=11)
ax.legend(fontsize=8.5,loc="upper left"); ax.spines[["top","right"]].set_visible(False)
plt.tight_layout(); plt.savefig(f"{OUT}/bce_collapse_en.png",bbox_inches="tight"); plt.close()

# multi-turn chat table
rows=[["1","DETECT prompt (as trained)","<ref>liver tumor</ref><box>[[438,211,469,241],…]</box>\n→ draws multiple boxes"],
      ["5",'JSON, exactly 3 keys','{"lesion_seen":"Yes",\n "confidence_caveat":"…",\n "next_step":"…"}'],
      ["7","2-line Vietnamese poem\nabout rain","Mưa rơi tí tách trên mái hiên,\nGột rửa phố phường, lòng thêm dịu êm."],
      ["8","17 × 23 = ?","391  ✓"]]
fig,ax=plt.subplots(figsize=(9.5,3.2),dpi=150); ax.axis("off")
ax.set_title("Multi-turn chat after fine-tune (1 positive case)",fontweight="bold",color=NAVY_H,fontsize=12,pad=12)
tbl=ax.table(cellText=rows,colLabels=["Turn","Request","Model reply"],colWidths=[0.08,0.30,0.62],cellLoc="left",loc="center")
tbl.auto_set_font_size(False); tbl.set_fontsize(9.5); tbl.scale(1,2.7)
for (r,c),cell in tbl.get_celld().items():
    cell.set_edgecolor("#d0d7de")
    if r==0: cell.set_facecolor(NAVY_H); cell.set_text_props(color="white",fontweight="bold")
    else:
        cell.set_facecolor("#f4f7fa" if r%2 else "#ffffff")
        if c==0: cell.set_text_props(fontweight="bold",color=NAVY_H)
plt.savefig(f"{OUT}/chat_test_multiturn_en.png",bbox_inches="tight"); plt.close()
print("EN charts done")

# ================= DECK (append to English paper deck) =================
NAVY=RGBColor(0x1A,0x3A,0x5C); SUBTXT=RGBColor(0xB8,0xCF,0xE0); FOOTBG=RGBColor(0xEE,0xF4,0xFB)
FOOTTX=RGBColor(0x6B,0x7F,0x95); CARD2=RGBColor(0xE8,0xF0,0xF8); GREY=RGBColor(0x75,0x70,0x70)
INK=RGBColor(0x1C,0x27,0x36); INFO=RGBColor(0x1F,0x55,0x8C); CAV=RGBColor(0xA0,0x46,0x46)
CAVBG=RGBColor(0xFB,0xF1,0xF1); PAGECL=RGBColor(0x5A,0x5A,0x5A); WHITE=RGBColor(0xFF,0xFF,0xFF)
TFONT,BFONT="Cambria","Calibri"; SW,SH=10.0,5.625
import re
SRC=r"C:/Users/loocn/Downloads/MedRegA_slide.pptx"
prs=Presentation(SRC); BLANK=prs.slide_layouts[0]

def _set(run,size,color,bold=False,italic=False,font=BFONT):
    run.font.size=Pt(size); run.font.bold=bold; run.font.italic=italic; run.font.name=font; run.font.color.rgb=color
def rect(slide,l,t,w,h,fill):
    sp=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False; return sp
def _runs(p,text):
    out=[]
    for seg in re.split(r"(\*\*.*?\*\*)",text):
        if not seg: continue
        r=p.add_run()
        if seg.startswith("**") and seg.endswith("**"): r.text=seg[2:-2]; out.append((r,True))
        else: r.text=seg; out.append((r,False))
    return out
def textbox(slide,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb=slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1); return tb,tf
def header(slide,title,subtitle=None):
    h=1.30 if subtitle else 0.97; rect(slide,0,0,SW,h,NAVY)
    tb,tf=textbox(slide,0.35,0.12 if subtitle else 0.20,9.3,0.60)
    for r,_ in _runs(tf.paragraphs[0],title): _set(r,20,WHITE,bold=True,font=TFONT)
    if subtitle:
        tb2,tf2=textbox(slide,0.37,0.74,9.2,0.45)
        r=tf2.paragraphs[0].add_run(); r.text=subtitle; _set(r,11,SUBTXT,italic=True)
    return h+0.12
def footer(slide,page=None):
    page=len(prs.slides._sldIdLst); rect(slide,0,5.34,SW,0.29,FOOTBG)
    tb,tf=textbox(slide,0.35,5.40,6.8,0.18); r=tf.paragraphs[0].add_run()
    r.text="MedRegA Improvements · Gemma 4 E4B (liver tumor, LiTS) · Deep Learning Seminar"; _set(r,8.5,FOOTTX)
    tb2,tf2=textbox(slide,9.0,5.40,0.8,0.18); tf2.paragraphs[0].alignment=PP_ALIGN.RIGHT
    r2=tf2.paragraphs[0].add_run(); r2.text=str(page); _set(r2,9,PAGECL)
def bullets(slide,items,l,t,w,h,base=12.5,gap=6):
    tb,tf=textbox(slide,l,t,w,h); first=True
    for text,lvl in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level=lvl; p.space_after=Pt(gap); p.space_before=Pt(0)
        pr=p.add_run(); pr.text=("▪  " if lvl==0 else "–  "); _set(pr,base,NAVY if lvl==0 else GREY,bold=(lvl==0))
        for r,isb in _runs(p,text): _set(r,base if lvl==0 else base-0.5,INK if lvl==0 else GREY,bold=isb)
    return tf
def notebox(slide,label,text,l,t,w,h,kind="info"):
    bg=CARD2 if kind=="info" else CAVBG; titlecl=INFO if kind=="info" else CAV
    rect(slide,l,t,w,h,bg); tb,tf=textbox(slide,l+0.12,t+0.07,w-0.24,h-0.14)
    r=tf.paragraphs[0].add_run(); r.text=label; _set(r,10,titlecl,bold=True)
    p2=tf.add_paragraph(); p2.space_before=Pt(1)
    for r,isb in _runs(p2,text): _set(r,9.5,INK,bold=isb)
def image(slide,path,l,t,w=None,h=None,caption=None):
    kw={}
    if w: kw["width"]=Inches(w)
    if h: kw["height"]=Inches(h)
    pic=slide.shapes.add_picture(path,Inches(l),Inches(t),**kw)
    if caption:
        tb,tf=textbox(slide,l,t+pic.height/914400+0.02,max(pic.width/914400,2.0),0.35)
        pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER; r=pp.add_run(); r.text=caption; _set(r,7.5,FOOTTX,italic=True)
    return pic
def table(slide,headers,rows,l,t,w,colw=None,fs=9,hfs=9,first_bold=True,rh=0.30):
    nr,nc=len(rows)+1,len(headers)
    g=slide.shapes.add_table(nr,nc,Inches(l),Inches(t),Inches(w),Inches(rh*nr)).table
    g.first_row=False; g.horz_banding=False
    if colw:
        for i,cw in enumerate(colw): g.columns[i].width=Inches(cw)
    for j,h in enumerate(headers):
        c=g.cell(0,j); c.fill.solid(); c.fill.fore_color.rgb=NAVY
        c.margin_left=Pt(4); c.margin_right=Pt(4); c.margin_top=Pt(2); c.margin_bottom=Pt(2); c.vertical_anchor=MSO_ANCHOR.MIDDLE
        for r,isb in _runs(c.text_frame.paragraphs[0],h): _set(r,hfs,WHITE,bold=True)
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=g.cell(i,j); c.fill.solid(); c.fill.fore_color.rgb=WHITE if i%2 else CARD2
            c.margin_left=Pt(4); c.margin_right=Pt(4); c.margin_top=Pt(1); c.margin_bottom=Pt(1); c.vertical_anchor=MSO_ANCHOR.MIDDLE
            for r,isb in _runs(c.text_frame.paragraphs[0],val): _set(r,fs,INK,bold=(isb or (first_bold and j==0)))
    return g
def notes(slide,text): slide.notes_slide.notes_text_frame.text=text
def newslide(): return prs.slides.add_slide(BLANK)

# ---- TITLE ----
s=newslide()
rect(s,0,0,SW,0.24,NAVY); rect(s,0,1.55,SW,2.30,NAVY)
tb,tf=textbox(s,0.5,1.75,9.0,1.0,MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text="MedRegA Improvements"; _set(r,38,WHITE,bold=True,font=TFONT)
tb2,tf2=textbox(s,0.5,2.80,9.0,0.9)
for i,line in enumerate(["Fine-tuning Gemma 4 E4B for liver-tumor detection on CT (LiTS)",
                         "Going deep on a small model — honest about evaluation & reliability"]):
    pp=tf2.paragraphs[0] if i==0 else tf2.add_paragraph(); pp.alignment=PP_ALIGN.CENTER
    rr=pp.add_run(); rr.text=line; _set(rr,15,SUBTXT)
rect(s,0,5.34,SW,0.29,FOOTBG); tb3,tf3=textbox(s,0.35,5.40,9,0.18)
r3=tf3.paragraphs[0].add_run(); r3.text="Deep Learning Seminar · 2026 · based on MedRegA (Wang et al., ICLR 2025)"; _set(r3,9,FOOTTX)
notes(s,"Opening. We take the opposite route to MedRegA: a small model, one disease (liver tumor), but going deep on evaluation methodology and reliability.")

# ---- A1 ----
s=newslide(); bt=header(s,"Clinical problem: why “knowing the region” matters","Part A — MedRegA (original method)")
bullets(s,[("Radiologists read scans as: view whole image → **localize region** → describe → diagnose with location.",0),
  ("A **region-agnostic** model compresses the whole image into one vector then generates text → **skips localization**.",0),
  ("Consequence: **wrong location leads to wrong diagnosis**, and the clinician cannot verify the answer.",1),
  ("MedRegA fixes exactly this: force the model to **emit region coordinates** so each statement anchors to pixels.",0)],0.4,bt,5.4,3.5)
image(s,f"{OUT}/fig2_left.png",5.85,2.05,w=4.05,caption="MedDr wrong region vs MedRegA correct region — Fig.2a")
footer(s)
notes(s,"Wrong region is not a minor wording error — it changes the diagnosis. That is why 'region-centric' exists.")

# ---- A2 ----
s=newslide(); bt=header(s,"MedRegA architecture (4 blocks)","Part A — original method")
bullets(s,[("Built on **InternVL 1.2 (~40B params)** — 4 blocks in series:",0),
  ("**Vision Encoder** (“eyes”): InternViT-6B, 448×448 → 32×32 token grid.",1),
  ("**Pixel Shuffle**: compress image tokens, keep 2D structure.",1),
  ("**MLP Connector** (“bridge”): project image features into the LLM space.",1),
  ("**LLM** (“brain”): Nous-Hermes-2-Yi-34B, text + reasoning.",1),
  ("**Box written as TEXT:** <ref>liver tumor</ref><box>[x1,y1,x2,y2]</box>, coords normalized **0–1000** → turns “localization” into “text generation” (no separate detection head).",0)],0.4,bt,5.5,3.6)
image(s,f"{OUT}/fig4_arch.png",5.9,2.5,w=3.95,caption="Vision Encoder→Alignment→LLM + box=token (Fig.4)")
footer(s)
notes(s,"Instead of a detection head, MedRegA teaches the model to WRITE coordinates as text — reusing the LLM's language power.")

# ---- A3 ----
s=newslide(); bt=header(s,"MedRegA training (2 stages) + 3 tasks","Part A — original method")
bullets(s,[("**Stage 1 — Alignment:** train the Connector to “translate” image↔text on image–caption pairs, NO boxes; Vision + LLM frozen.",0),
  ("**Stage 2 — Instruction tuning:** on MedRegInstruct (**790K samples** = 550K region–text + 240K grounded report, 8 modalities, EN–CN bilingual); tune the LLM.",0),
  ("**3 region-centric tasks:** Region→Text · **Text→Region (our project does this task)** · Grounded Report.",0),
  ("Loss = **plain language cross-entropy** (no geometric loss for boxes); 3D images use only **1 central slice**; **16× H800, several days**.",0)],0.4,bt,9.2,2.2)
image(s,f"{OUT}/fig4_tasks.png",1.35,3.55,w=7.3,caption="3 region-centric tasks: Region→Text · Text→Region · Grounded Report (Fig.4)")
footer(s)
notes(s,"Two stages: first teach the 'bridge' to link image–text, then teach the 'brain' to follow instructions. Boxes are just text, so plain language loss suffices.")

# ---- A4 ----
s=newslide(); bt=header(s,"Limitations of the original paper (our openings)","Part A — original method")
bullets(s,[("Most of the load sits on the **40B LLM**; the “eyes” & “bridge” are nearly frozen → localization “works but is not sharp” (low IoU, misses regions in multi-focal cases).",0),
  ("**No abstain** mechanism (refusing when unsure); **no negative cases** → cannot measure whether the model invents boxes.",0),
  ("Evaluation is **per-image**; many slices from one patient are counted as independent samples (**pseudoreplication**); no confidence intervals.",0),
  ("Cost is **out of reach** for a student project.",0)],0.4,bt,5.5,3.4)
image(s,f"{OUT}/fig2_right.png",6.3,1.65,w=3.35,caption="w/ vs w/o Region: regions contribute a lot (Fig.2b)")
footer(s)
notes(s,"The paper is strong on scale and benchmarks. We don't claim it lacks region metrics — we add a clinically-closer evaluation layer.")

# ---- A5 (map) ----
s=newslide(); bt=header(s,"Gap & improvement map","Bridge A → B: 4 limitations ↔ 4 improvements")
table(s,["Original limitation","Our improvement (→Slide)"],
  [["Frozen eyes/bridge (localization not sharp)","Unlock vision (→S6)"],
   ["1 central slice · no negatives · coarse boxes","Multi-slice + negatives + multi-box (→S7)"],
   ["Cost 16× H800 out of reach","LoRA, 1 GPU (→S8)"],
   ["Per-image · no CI · no abstain","Per-patient + CI + selective prediction (→S9)"]],
  0.7,bt+0.25,8.6,colw=[4.3,4.3],fs=11,hfs=11,rh=0.55)
footer(s)
notes(s,"The bridge: 4 limitations on the left, 4 fixes on the right — unlock vision, multi-slice+negatives data, 1-GPU training, and a clinically-closer evaluation layer.")

# ---- B1 backbone ----
s=newslide(); bt=header(s,"Improvement (1): swap backbone — 5× smaller, 1 GPU","Part B — our improvements")
table(s,["","MedRegA (orig.)","Ours"],
  [["Backbone","InternVL 1.2 (~40B)","Gemma 4 E4B (~8B)"],
   ["Vision","InternViT-6B (frozen)","SigLIP — WE UNLOCK & fine-tune"],
   ["Runs on","16× H800","1×A100 / Colab"]],0.4,bt+0.05,5.5,colw=[1.2,2.3,2.0],fs=9.5,hfs=9.5,rh=0.42)
bullets(s,[("**Unlocking the vision tower is the key step**: frozen → pIoU 0.015; unlocked → 0.27.",0),
  ("⚠️ 1 config for illustration, **not multi-seed** — shows the DIRECTION only.",1),
  ("⚠️ Do NOT say “Gemma beats MedRegA” — it must be fine-tuned on medical data to compete.",1)],0.4,bt+2.05,5.5,1.6,base=11)
image(s,f"{OUT}/iou_frozen_vs_unfrozen_en.png",6.2,bt+0.1,w=3.5)
footer(s)
notes(s,"First contribution: a feasibility proof — reproducing region-detection on a much smaller model, on a single GPU. Philosophy: trade 'breadth' (8 modalities) for 'depth' (1 task).")

# ---- B2 data ----
s=newslide(); bt=header(s,"Improvement (2): data — multi-slice, negatives, multi-box","Part B — our improvements")
bullets(s,[("**Multi-slice** (vs paper's 1 slice): per patient, 2 largest-tumor slices + 3 random.",0),
  ("**Negative cases (no tumor):** teach the model “when NOT to draw a box” — missing in the paper.",0),
  ("**Multi-box** (connected-components: split touching tumor-pixel clusters): 1 box/lesion, drop <30px. **35.5% of positive slices have ≥2 lesions**.",0),
  ("**Liver-CT-appropriate augmentation**: shift/small rotate ±12°/scale — NO flips (liver has anatomical orientation).",0),
  ("**Dataset:** 1211 samples (558 pos / 653 neg), **split by patient**.",0)],0.4,bt,5.2,3.4,base=11.5)
image(s,f"{OUT}/data/data_liver/images/liver_002_pos_z457.png",5.75,2.35,w=2.0,caption="WITH tumor")
image(s,f"{OUT}/data/data_liver/images/liver_002_neg_z389.png",7.9,2.35,w=2.0,caption="NO tumor")
footer(s)
notes(s,"The paper uses only the central slice and no negatives. We add both — but 5 slices of one patient are highly correlated, so we count by patient, not by slice.")

# ---- B3 training ----
s=newslide(); bt=header(s,"Improvement (3): training setup","Part B — our improvements")
bullets(s,[
  ("**LoRA on the LLM** (small inserted matrices, keeps base weights): learn only a small delta, don't damage base skills.",0),
  ("**Full fine-tune of the Vision tower** (fp32, LR 1e-5) — unlock the “eyes” to learn medical features (the key step, see previous slide).",0),
  ("**Two separate learning-rate groups** (LoRA 2e-4 / vision 1e-5) for stable training.",0),
  ("**Loss:** standard language cross-entropy — box coordinates are text tokens, no separate geometric loss.",0),
  ("Runs in **bf16** on a single GPU — fits a student budget.",0)],0.4,bt,9.2,3.4,base=12.5)
footer(s)
notes(s,"Compact training setup: LoRA keeps base LLM skills, unlocking vision learns medical features, two LR groups keep it stable. Boxes are text tokens, so plain language loss suffices; runs on 1 GPU.")

# ---- B4 eval ----
s=newslide(); bt=header(s,"Improvement (4): stricter EVALUATION — our strongest contribution","Part B — our improvements")
bullets(s,[("**Per-patient + bootstrap CI** instead of per-image → fixes pseudoreplication; true n ~25 positive patients → **wide CI, stated honestly**.",0),
  ("**Separate DETECTION from LOCALIZATION:** detection = a box overlapping GT; false-positive = a box on a negative case; localization = IoU reported separately.",0),
  ("**Selective prediction:** two free signals — **coordinate-token logprob** + **spatial consistency** → risk–coverage + conformal threshold → **triage** (flag for a doctor).",0)],0.4,bt,5.5,3.2,base=11.5)
notebox(s,"⚠ Conformal caveat","PoC — with small n (~25) the statistical guarantee is loose; needs a larger calibration set.",0.4,bt+3.05,5.5,0.7,kind="caveat")
image(s,f"{OUT}/eval/risk_coverage.png",6.15,bt+0.15,w=3.5,caption="risk–coverage: dropping low-confidence cases raises IoU")
footer(s)
notes(s,"Our biggest contribution is the evaluation method — is there a tumor, where is it, and when to defer to a doctor — all read straight from the generation, no extra network.")

# ---- C1 detection ----
s=newslide(); bt=header(s,"Result: Detection (tumor / no tumor) — GOOD","Part C — Results (our model, n=25 pos / 2 neg)")
table(s,["Metric","Zero-shot","Fine-tuned"],
  [["Detection F1 (patient)","0.33","0.89"],
   ["Sensitivity (catch tumor)","20%","80% (20/25)"],
   ["False-positive (neg. slice)","2.2%","1.5% (2/135)"],
   ["Specificity/Prec (patient)","—","100% (2/2 — n=2!)"]],0.4,bt+0.05,5.5,colw=[2.7,1.3,1.5],fs=9.5,hfs=9.5,rh=0.40)
notebox(s,"⚠ Mandatory caveat","Spec/Prec 100% rests on ONLY 2 tumor-free patients → statistically very weak; it is within-patient slice discrimination, NOT clinical specificity.",0.4,bt+2.05,5.5,0.95,kind="caveat")
image(s,f"{OUT}/detect_before_after_en.png",6.1,bt+0.2,w=3.6,caption="Before → after fine-tune")
footer(s)
notes(s,"We catch 20/25 patients with low FP. But that 100% rests on 2 people — we do not use it to claim clinical specificity. Zero-shot almost always says 'No liver tumor is found' → fine-tuning teaches the task from zero.")

# ---- C2 localization ----
s=newslide(); bt=header(s,"Result: Localization + why the model MISSES tumors","Part C — Results")
bullets(s,[("**Localization is MODERATE** (separated): per-patient **pIoU ≈ 0.27**, CI95 [0.19, 0.36]; recall@0.25 ≈ **54%** (small 35% / medium 54% / large 74%).",0),
  ("**Misses are NOT only small tumors:** per-lesion, missed lesions are ~2× smaller AND lower-contrast; but per-patient the sizes are nearly equal.",0),
  ("**Low-contrast/isointense** (tumor as bright as liver) is what discriminates even large tumors — pid116, a 5.63% tumor, was still missed.",1),
  ("Suspect the wide CT window (WL40/WW400) → a **narrow liver window** is the next thing to try.",1)],0.4,bt,5.5,3.4,base=11)
image(s,f"{OUT}/pid116_contrast.png",6.15,bt+0.05,w=3.5,caption="pid116: 5.63% tumor still hard to see (isointense)")
image(s,f"{OUT}/missed_patients_v3.png",6.15,bt+2.35,w=3.5,caption="overview of missed cases")
footer(s)
notes(s,"Localization is only moderate, we don't hide it. Misses lean toward small AND low-contrast tumors — both contribute. Note: the two right-side medical figures still carry Vietnamese sublabels (generated earlier).")

# ---- C3 no-forget ----
s=newslide(); bt=header(s,"Result: fine-tuning does NOT cause “forgetting”","Part C — Results")
bullets(s,[("Concern: does narrow fine-tuning destroy conversational ability? → tested with **multi-turn chat**, saved as JSON.",0),
  ("The model both **draws boxes** and: returns exact-key JSON, writes a poem, computes 17×23=**391**, explains bilingually.",0),
  ("**Why preserved:** LoRA keeps base weights → adds skills without erasing old ones.",0),
  ("⚠️ This is a **QUALITATIVE** check (a few cases), not a full general benchmark.",1)],0.4,bt,4.6,3.4,base=11.5)
image(s,f"{OUT}/chat_test_multiturn_en.png",5.15,bt+0.35,w=4.55,caption="multi-turn chat (real data)")
footer(s)
notes(s,"The model draws the tumor AND follows instructions precisely, explains bilingually, writes poetry, computes 17×23=391. Thanks to LoRA keeping base weights — new skills without forgetting. Qualitative check; no MMLU measured, so we don't claim full preservation.")

# ---- D1 compare ----
s=newslide(); bt=header(s,"Comparison with the original paper","Part D — NOT a direct number match; comparing properties")
table(s,["Aspect","MedRegA (orig.)","Ours"],
  [["Scale","~40B, 16×H800","~8B, 1×A100"],
   ["Into 3D","1 central slice","multi-slice + negatives"],
   ["Region labels","coarse merged box","multi-box"],
   ["Eval unit","per-image","per-patient + CI"],
   ["Detect vs Localize","merged","separated"],
   ["Negatives / FP · Uncertainty","none","FP + selective prediction"],
   ["Grounded report","yes","not yet (→Future)"]],0.4,bt,5.6,colw=[1.8,1.9,1.9],fs=8.8,hfs=9,rh=0.33)
notebox(s,"Localization — read it correctly","Orig. IoU ~0.23 vs ours pIoU ~0.27 (same 0–1 scale) — different data/metric, NOT directly comparable, NO better/worse claim. pIoU is PENALIZED, so it's stricter.",6.15,bt+0.1,3.5,1.5,kind="info")
notebox(s,"The real contribution","Not the score — but an honest evaluation + reliability layer added on top of the paper's benchmark.",6.15,bt+1.75,3.5,1.1,kind="info")
footer(s)
notes(s,"We don't claim to beat MedRegA. The contribution is turning a localization benchmark into a clinically-closer evaluation workflow: is there a tumor, where, and when not to trust it.")

# ---- D2 future1 ----
s=newslide(); bt=header(s,"Future work (1): richer images + more trustworthy","Part D — [MEASURED] ties to Slides C1–C2 · [ASSUMED] from design")
table(s,["#","Limitation (source)","Next direction","Cost"],
  [["①","Miss isointense tumors [MEASURED]","Multi-window (narrow liver window)","🟢 very low"],
   ["②","3D→2D info loss [ASSUMED]","2.5D stack (z−1/z/z+1 into RGB)","🟢 low"],
   ["③","LiTS single-phase [ASSUMED]","Multi-phase CT — a DATA ceiling, not the model","🔴 needs data"],
   ["⑤","2 negatives, 1 dataset [MEASURED]","External validation + FROC + calibration","🟡 moderate"]],0.4,bt+0.15,9.2,colw=[0.4,3.4,3.9,1.5],fs=10,hfs=10,rh=0.52)
footer(s)
notes(s,"Track 1 — richer images (multi-window, 2.5D, multi-phase) and more trustworthy results (external validation, more negatives). ①⑤ tie straight to results; ②③ are design-based, so marked not-yet-measured. External validation = testing on other CT sets (IRCADb/MSD). FROC = multi-lesion detection curve. calibration = aligning confidence with true probability.")

# ---- D3 future2 ----
s=newslide(); bt=header(s,"Future work (2): grounded reporting · generalist","Part D — [MEASURED] · [SCOPE] · [ASSUMED]")
table(s,["#","Limitation (source)","Next direction","Cost"],
  [["④","pIoU only 0.27 [MEASURED]","MedSAM tight masks + polygon-as-token","🟡 moderate"],
   ["⑥","Detection only [SCOPE]","Grounded report (each sentence anchored to a box)","🟡 moderate"],
   ["⑦","Liver tumor only [ASSUMED]","Multi-organ generalist","🔴 long-term"],
   ["⑧","Token-only CE for coordinates [SCOPE]","Smarter loss balancing (uncertainty weighting + geometric GIoU)","🟢 low"]],0.4,bt+0.15,9.2,colw=[0.4,3.2,4.1,1.5],fs=10,hfs=10,rh=0.52)
footer(s)
notes(s,"Track 2 — from pointing to reporting-with-evidence, then a multi-organ generalist; and smarter loss (geometric GIoU) instead of token-only coordinate CE. Verify tool names (MedSAM/FROC/IRCADb) before submission.")

# ---- D4 conclusion ----
s=newslide(); bt=header(s,"Conclusion","Part D")
bullets(s,[("**Proof of a mechanism:** region-centric (localize-then-read) is **feasible at ~8B, 1 GPU**.",0),
  ("**Core result:** Detection F1 **0.89** · pIoU **0.27** · FP **1.5%** — but the value is the **honest evaluation + selective prediction**, not the score.",0),
  ("**Honest:** PoC on 1 dataset / 1 task; small n, wide CI; no clinical-system claim.",0),
  ("**Expansion map:** richer images → more trustworthy → grounded reporting & generalist.",0)],0.4,bt,9.2,3.2,base=12.5)
notebox(s,"Closing line","We don't claim to have arrived — we map the road and say clearly where it's still steep.",0.4,bt+3.05,9.2,0.6,kind="info")
footer(s)
notes(s,"We don't claim to beat MedRegA on scale or SOTA. The main contribution is turning a region-localization benchmark into a clinically-closer, honest evaluation workflow.")

out="seminar/cai_tien/seminar_full_EN_2026-07-05.pptx"
prs.save(out)
print("SAVED",out,"| n_slides =",len(prs.slides._sldIdLst))
