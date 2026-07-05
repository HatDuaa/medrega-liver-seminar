# -*- coding: utf-8 -*-
"""Dựng slides_cai_tien.pptx theo theme deck paper MedRegA (navy). Chạy cwd = repo root."""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY   = RGBColor(0x1A,0x3A,0x5C)
SUBTXT = RGBColor(0xB8,0xCF,0xE0)
FOOTBG = RGBColor(0xEE,0xF4,0xFB)
FOOTTX = RGBColor(0x6B,0x7F,0x95)
CARD   = RGBColor(0xF7,0xF9,0xFC)
CARD2  = RGBColor(0xE8,0xF0,0xF8)
GREY   = RGBColor(0x75,0x70,0x70)
INK    = RGBColor(0x1C,0x27,0x36)
INFO   = RGBColor(0x1F,0x55,0x8C)
CAV    = RGBColor(0xA0,0x46,0x46)
CAVBG  = RGBColor(0xFB,0xF1,0xF1)
PAGECL = RGBColor(0x5A,0x5A,0x5A)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
GREEN  = RGBColor(0x2E,0x7D,0x32)
TFONT, BFONT = "Cambria", "Calibri"
SW, SH = 10.0, 5.625

SRC_DECK = r"C:/Users/loocn/Downloads/MedRegA_slide.pptx"   # deck paper gốc (12 slide)
prs = Presentation(SRC_DECK)                                # mở gốc → kế thừa master/theme/font
BLANK = prs.slide_layouts[0]                                # layout 'DEFAULT' rỗng của deck gốc
# slide của mình sẽ được APPEND sau 12 slide paper -> 1 deck chung

def _set(run, size, color, bold=False, italic=False, font=BFONT):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = color

def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def _runs(p, text):
    """Split **bold** markup into runs."""
    import re
    parts = re.split(r"(\*\*.*?\*\*)", text)
    out = []
    for seg in parts:
        if not seg: continue
        r = p.add_run()
        if seg.startswith("**") and seg.endswith("**"):
            r.text = seg[2:-2]; out.append((r, True))
        else:
            r.text = seg; out.append((r, False))
    return out

def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tb, tf

def header(slide, title, subtitle=None):
    h = 1.30 if subtitle else 0.97
    rect(slide, 0, 0, SW, h, NAVY)
    tb, tf = textbox(slide, 0.35, 0.12 if subtitle else 0.20, 9.3, 0.60, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    for r,_ in _runs(p, title): _set(r, 20, WHITE, bold=True, font=TFONT)
    if subtitle:
        tb2, tf2 = textbox(slide, 0.37, 0.74, 9.2, 0.45)
        p2 = tf2.paragraphs[0]
        r = p2.add_run(); r.text = subtitle; _set(r, 11, SUBTXT, italic=True)
    return h + 0.12

def footer(slide, page=None):
    page = len(prs.slides._sldIdLst)   # số trang THẬT theo vị trí trong deck chung
    rect(slide, 0, 5.34, SW, 0.29, FOOTBG)
    tb, tf = textbox(slide, 0.35, 5.40, 6.5, 0.18)
    r = tf.paragraphs[0].add_run(); r.text = "Cải tiến MedRegA · Gemma 4 E4B (u gan LiTS) · Seminar Học sâu"
    _set(r, 8.5, FOOTTX)
    tb2, tf2 = textbox(slide, 9.0, 5.40, 0.8, 0.18)
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r2 = tf2.paragraphs[0].add_run(); r2.text = str(page); _set(r2, 9, PAGECL)

def bullets(slide, items, l, t, w, h, base=12.5, gap=6):
    """items: list of (text, level). level 0/1. Supports **bold**."""
    tb, tf = textbox(slide, l, t, w, h)
    first = True
    for text, lvl in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap); p.space_before = Pt(0)
        bullet_char = "▪" if lvl == 0 else "–"
        pr = p.add_run(); pr.text = bullet_char + "  "
        _set(pr, base, NAVY if lvl == 0 else GREY, bold=(lvl==0))
        for r, isb in _runs(p, text):
            _set(r, base if lvl==0 else base-0.5, INK if lvl==0 else GREY, bold=isb)
    return tf

def notebox(slide, label, text, l, t, w, h, kind="info"):
    bg = CARD2 if kind=="info" else CAVBG
    titlecl = INFO if kind=="info" else CAV
    rect(slide, l, t, w, h, bg)
    tb, tf = textbox(slide, l+0.12, t+0.07, w-0.24, h-0.14)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = label; _set(r, 10, titlecl, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(1)
    for r, isb in _runs(p2, text): _set(r, 9.5, INK, bold=isb)

def image(slide, path, l, t, w=None, h=None, caption=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    pic = slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)
    if caption:
        cw = (pic.width/914400)
        tb, tf = textbox(slide, l, t + pic.height/914400 + 0.02, max(cw, 2.0), 0.35)
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = caption; _set(r, 7.5, FOOTTX, italic=True)
    return pic

def table(slide, headers, rows, l, t, w, colw=None, fs=9, hfs=9, first_bold=True, rh=0.30):
    nr, nc = len(rows)+1, len(headers)
    gtbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(t), Inches(w), Inches(rh*nr)).table
    gtbl.first_row = False; gtbl.horz_banding = False
    if colw:
        for i, cw in enumerate(colw): gtbl.columns[i].width = Inches(cw)
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.margin_left=Pt(4); c.margin_right=Pt(4); c.margin_top=Pt(2); c.margin_bottom=Pt(2)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = c.text_frame.paragraphs[0]
        for r, isb in _runs(para, htxt): _set(r, hfs, WHITE, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gtbl.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else CARD2
            c.margin_left=Pt(4); c.margin_right=Pt(4); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = c.text_frame.paragraphs[0]
            for r, isb in _runs(para, val):
                _set(r, fs, INK, bold=(isb or (first_bold and j==0)))
    return gtbl

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def newslide():
    return prs.slides.add_slide(BLANK)

# =========================================================
# SLIDE 0 — TITLE
# =========================================================
s = newslide()
rect(s, 0, 0, SW, 0.24, NAVY)
rect(s, 0, 1.55, SW, 2.30, NAVY)
tb, tf = textbox(s, 0.5, 1.75, 9.0, 1.0, MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Cải tiến MedRegA"; _set(r, 40, WHITE, bold=True, font=TFONT)
tb2, tf2 = textbox(s, 0.5, 2.80, 9.0, 0.9, MSO_ANCHOR.TOP)
for line in ["Fine-tune Gemma 4 E4B phát hiện u gan trên CT (LiTS)",
             "Đi sâu trên mô hình nhỏ — trung thực về đánh giá & độ tin cậy"]:
    pp = tf2.add_paragraph() if line != "Fine-tune Gemma 4 E4B phát hiện u gan trên CT (LiTS)" else tf2.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = line; _set(rr, 15, SUBTXT)
rect(s, 0, 5.34, SW, 0.29, FOOTBG)
tb3, tf3 = textbox(s, 0.35, 5.40, 9, 0.18)
r3 = tf3.paragraphs[0].add_run(); r3.text = "Seminar Học sâu · 2026 · dựa trên MedRegA (Wang et al., ICLR 2025)"
_set(r3, 9, FOOTTX)
notes(s, "Slide mở. Chúng em chọn hướng ngược với MedRegA: mô hình nhỏ, một bệnh (u gan), nhưng đào sâu phương pháp đánh giá và độ tin cậy.")

# =========================================================
# PHẦN A — PHƯƠNG PHÁP GỐC
# =========================================================
# S1
s = newslide(); bt = header(s, "Vấn đề lâm sàng: vì sao cần “biết vùng”", "Phần A — Phương pháp gốc MedRegA")
bullets(s, [
  ("Bác sĩ đọc phim: nhìn toàn ảnh → **khoanh vùng** → mô tả → chẩn đoán kèm vị trí.", 0),
  ("Model **region-agnostic** (không biết vùng): nén cả ảnh thành 1 vector rồi sinh chữ → bỏ mất bước khoanh vùng.", 0),
  ("Hệ quả: **sai vị trí kéo theo sai chẩn đoán**, bác sĩ không kiểm chứng được.", 1),
  ("MedRegA vá đúng chỗ này: bắt model **sinh kèm toạ độ vùng** → mỗi nhận định neo vào đúng pixel.", 0),
], 0.4, bt, 5.4, 3.5)
image(s, "seminar/code/fig2_left.png", 5.85, 2.05, w=4.05,
      caption="MedDr sai vùng vs MedRegA đúng vùng — Fig.2a")
footer(s, 2)
notes(s, "Câu nói: Sai vùng không phải lỗi nhỏ về câu chữ — nó đổi luôn chẩn đoán. Đó là lý do 'region-centric' ra đời.\n\nregion-agnostic = mô hình không có khái niệm vùng, xử lý cả ảnh như một khối.")

# S2
s = newslide(); bt = header(s, "Kiến trúc MedRegA (4 khối)", "Phần A — Phương pháp gốc")
bullets(s, [
  ("Xây trên **InternVL 1.2 (~40B tham số)** — 4 khối nối tiếp:", 0),
  ("**Vision Encoder** (“mắt”): InternViT-6B, ảnh 448×448 → lưới token 32×32.", 1),
  ("**Pixel Shuffle**: dồn/nén token ảnh, giữ cấu trúc 2D.", 1),
  ("**MLP Connector** (“cầu”): chiếu đặc trưng ảnh sang không gian LLM.", 1),
  ("**LLM** (“não”): Nous-Hermes-2-Yi-34B, sinh chữ + suy luận.", 1),
  ("**Box viết thành CHỮ:** <ref>liver tumor</ref><box>[x1,y1,x2,y2]</box>, toạ độ chuẩn hoá **0–1000** → biến “khoanh vùng” thành “sinh chữ” (không cần đầu phát hiện riêng).", 0),
], 0.4, bt, 5.5, 3.6)
image(s, "seminar/code/fig4_arch.png", 5.9, 2.5, w=3.95,
      caption="Luồng Vision Encoder→Alignment→LLM + box=token (Fig.4)")
footer(s, 3)
notes(s, "Câu nói: Thay vì gắn thêm 'đầu phát hiện', MedRegA dạy model VIẾT toạ độ ra như một câu chữ — tận dụng thẳng sức mạnh ngôn ngữ.")

# S3
s = newslide(); bt = header(s, "Huấn luyện MedRegA (2 giai đoạn) + 3 task", "Phần A — Phương pháp gốc")
bullets(s, [
  ("**GĐ1 — Alignment:** dạy Connector “phiên dịch” ảnh↔chữ trên cặp ảnh–caption, KHÔNG có box; Vision + LLM đóng băng.", 0),
  ("**GĐ2 — Instruction tuning:** trên MedRegInstruct (**790K mẫu** = 550K region–text + 240K grounded report, 8 modality, song ngữ Anh–Trung); chỉnh LLM.", 0),
  ("**3 task region-centric:** Region→Text · **Text→Region (đồ án làm task này)** · Grounded Report.", 0),
  ("Loss = **cross-entropy ngôn ngữ thuần** (không có loss hình học cho box); ảnh 3D chỉ lấy **1 lát trung tâm**; **16× H800, vài ngày**.", 0),
], 0.4, bt, 9.2, 2.2)
image(s, "seminar/code/fig4_tasks.png", 1.35, 3.55, w=7.3,
      caption="3 task region-centric: Region→Text · Text→Region · Grounded Report (Fig.4)")
footer(s, 4)
notes(s, "Câu nói: Hai giai đoạn — đầu tiên dạy 'cầu' nối ảnh–chữ, sau đó dạy 'não' làm việc theo lệnh. Box chỉ là chữ nên vẫn dùng loss ngôn ngữ bình thường.")

# S4
s = newslide(); bt = header(s, "Hạn chế của paper gốc (chỗ để mình lấp)", "Phần A — Phương pháp gốc")
bullets(s, [
  ("Dồn gánh nặng lên **LLM 40B**; “mắt” & “cầu” gần như bất động → định vị “được nhưng không sắc” (IoU thấp, sót vùng ở ca nhiều ổ).", 0),
  ("**Không có abstain** (từ chối khi không chắc); **không dùng ca âm** → không đo được model có bịa box.", 0),
  ("Đánh giá **per-image**, nhiều lát cùng bệnh nhân tính như mẫu độc lập (**pseudoreplication**); không khoảng tin cậy.", 0),
  ("Chi phí **ngoài tầm** một đồ án sinh viên.", 0),
], 0.4, bt, 5.5, 3.4)
image(s, "seminar/code/fig2_right.png", 6.3, 1.65, w=3.35,
      caption="w/ vs w/o Region: vùng đóng góp lớn (Fig.2b)")
footer(s, 5)
notes(s, "Câu nói: Paper mạnh về quy mô và benchmark. Chúng em không nói nó thiếu metric vùng — nó có 4 hạn chế để mình lấp.\n\nIoU = độ chồng lấn box (0–1). pseudoreplication = coi các lát phụ thuộc như độc lập → thổi phồng cỡ mẫu.")

# S4b
s = newslide(); bt = header(s, "Khoảng trống & bản đồ cải tiến", "Bản lề A → B: 4 hạn chế ↔ 4 cải tiến")
table(s, ["Hạn chế gốc", "Cải tiến của nhóm (→Slide)"], [
  ["Mắt/cầu bất động (định vị không sắc)", "Mở khoá vision (→S6)"],
  ["1 lát trung tâm · không ca âm · box gộp thô", "Đa lát + ca âm + multi-box (→S7)"],
  ["Chi phí 16× H800 ngoài tầm", "LoRA 1 GPU (→S8)"],
  ["Per-image · không CI · không abstain", "Per-patient + CI + selective prediction (→S9)"],
], 0.7, bt+0.25, 8.6, colw=[4.3,4.3], fs=11, hfs=11, rh=0.55)
footer(s, 6)
notes(s, "Câu nói: Bản lề — 4 hạn chế bên trái, 4 lớp vá bên phải: mở khoá thị giác, dữ liệu đa lát+ca âm, huấn luyện 1 GPU, và lớp đánh giá gần lâm sàng.")

# =========================================================
# PHẦN B — CẢI TIẾN
# =========================================================
# S5
s = newslide(); bt = header(s, "Cải tiến (1): đổi backbone — nhỏ 5×, chạy 1 GPU", "Phần B — Cải tiến của nhóm")
table(s, ["", "MedRegA gốc", "Của nhóm"], [
  ["Backbone", "InternVL 1.2 (~40B)", "Gemma 4 E4B (~8B)"],
  ["Vision", "InternViT-6B (đóng băng)", "SigLIP — MỞ KHOÁ fine-tune"],
  ["Chạy trên", "16× H800", "1×A100 / Colab"],
], 0.4, bt+0.05, 5.5, colw=[1.2,2.3,2.0], fs=9.5, hfs=9.5, rh=0.42)
bullets(s, [
  ("**Mở vision tower = bước tạo khác biệt**: đóng băng → pIoU 0.015; mở khoá → 0.27.", 0),
  ("⚠️ 1 cặp cấu hình minh hoạ, **chưa lặp seed** — chỉ cho thấy HƯỚNG.", 1),
  ("⚠️ KHÔNG nói “Gemma giỏi hơn MedRegA” — phải fine-tune y khoa mới cạnh tranh.", 1),
], 0.4, bt+2.05, 5.5, 1.6, base=11)
image(s, "seminar/code/iou_frozen_vs_unfrozen.png", 6.2, bt+0.25, w=3.6)
footer(s, 7)
notes(s, "Câu nói: Đóng góp đầu tiên là chứng minh khả thi — tái tạo năng lực khoanh vùng trên mô hình nhỏ hơn nhiều, chỉ 1 GPU. Triết lý: đổi 'rộng' (8 modality) lấy 'sâu' (1 task).")

# S6
s = newslide(); bt = header(s, "Cải tiến (2): dữ liệu — đa lát, ca âm, multi-box", "Phần B — Cải tiến")
bullets(s, [
  ("**Đa lát** (khác paper 1 lát): mỗi bệnh nhân 2 lát u lớn + 3 ngẫu nhiên.", 0),
  ("**Ca ÂM (không u):** dạy model “khi nào KHÔNG vẽ box” — paper thiếu.", 0),
  ("**Multi-box** (connected-components: tách cụm pixel u dính nhau): 1 box/ổ, lọc <30px. **35.5% lát dương có ≥2 ổ**.", 0),
  ("**Augment hợp CT gan**: dịch/xoay nhẹ ±12°/scale — KHÔNG lật (gan có hướng giải phẫu).", 0),
  ("**Bộ dữ liệu:** 1211 mẫu (558 dương/653 âm), **split theo bệnh nhân**.", 0),
], 0.4, bt, 5.2, 3.4, base=11.5)
image(s, "seminar/code/data/data_liver/images/liver_002_pos_z457.png", 5.75, 2.35, w=2.0, caption="CÓ u")
image(s, "seminar/code/data/data_liver/images/liver_002_neg_z389.png", 7.9, 2.35, w=2.0, caption="KHÔNG u")
footer(s, 7+1)
notes(s, "Câu nói: Paper chỉ dùng lát trung tâm, không có ca âm. Chúng em thêm cả hai — nhưng 5 lát của một người tương quan cao, nên đếm bằng bệnh nhân, không đếm bằng lát.\n\nconnected-components = gán nhãn từng vùng pixel liên thông (scipy.ndimage.label).")

# S7
s = newslide(); bt = header(s, "Cải tiến (3): thiết lập huấn luyện", "Phần B — Cải tiến")
bullets(s, [
  ("**LoRA trên LLM** (Low-Rank Adaptation — chèn ma trận nhỏ, giữ nguyên trọng số gốc): chỉ học thêm phần nhỏ, không phá năng lực gốc.", 0),
  ("**Full-finetune Vision tower** (fp32, LR 1e-5) — mở khoá “mắt” để học đặc trưng y khoa (bước tạo khác biệt, xem slide trước).", 0),
  ("**2 nhóm learning-rate riêng** (LoRA 2e-4 / vision 1e-5) để huấn luyện ổn định, không nổ số.", 0),
  ("**Loss:** cross-entropy ngôn ngữ tiêu chuẩn — toạ độ box coi như token chữ, không cần loss hình học riêng.", 0),
  ("Chạy **bf16**, chỉ **1 GPU** — phù hợp tài nguyên sinh viên.", 0),
], 0.4, bt, 9.2, 3.4, base=12.5)
footer(s)
notes(s, "Câu nói: Thiết lập huấn luyện gọn — LoRA cho LLM giữ năng lực gốc, mở khoá vision để học đặc trưng y khoa, 2 nhóm LR để ổn định. Box là token chữ nên dùng loss ngôn ngữ bình thường; chạy 1 GPU.")

# S8
s = newslide(); bt = header(s, "Cải tiến (4): ĐÁNH GIÁ chặt hơn — đóng góp mạnh nhất", "Phần B — Cải tiến")
bullets(s, [
  ("**Per-patient + bootstrap CI** thay per-image → sửa pseudoreplication; n thật ~25 dương → **CI rộng, nói thẳng**.", 0),
  ("**Tách PHÁT HIỆN khỏi ĐỊNH VỊ:** detect = box chồng GT; FP = ca âm bị vẽ box; localization = IoU riêng.", 0),
  ("**Selective prediction:** 2 tín hiệu miễn phí — **logprob toạ độ** + **nhất quán không gian** → risk–coverage + conformal → **triage** (gắn cờ cho bác sĩ).", 0),
], 0.4, bt, 5.5, 3.2, base=11.5)
notebox(s, "⚠ Caveat conformal", "PoC — với n nhỏ (~25) bảo đảm thống kê còn lỏng, cần calibration set lớn hơn.",
        0.4, bt+3.05, 5.5, 0.7, kind="caveat")
image(s, "seminar/code/eval/risk_coverage.png", 6.15, bt+0.15, w=3.5,
      caption="risk–coverage: bỏ ca kém tin → IoU tăng")
footer(s, 10)
notes(s, "Câu nói: Đóng góp lớn nhất ở phương pháp đánh giá — có/không u, khoanh ở đâu, khi nào đưa bác sĩ xem lại — rút thẳng từ lần sinh, không thêm mạng phụ.\n\nbootstrap CI = lấy mẫu lại nhiều lần ước lượng dao động. conformal = đặt ngưỡng có bảo đảm thống kê.")

# =========================================================
# PHẦN C — KẾT QUẢ
# =========================================================
# S9
s = newslide(); bt = header(s, "Kết quả: Phát hiện (có/không u) — TỐT", "Phần C — Kết quả (model của nhóm, n=25 dương / 2 âm)")
table(s, ["Chỉ số", "Zero-shot", "Fine-tune"], [
  ["Detection F1 (bệnh nhân)", "0.33", "0.89"],
  ["Sensitivity (bắt u)", "20%", "80% (20/25)"],
  ["False-positive (lát âm)", "2.2%", "1.5% (2/135)"],
  ["Specificity/Prec (bệnh nhân)", "—", "100% (2/2 — n=2!)"],
], 0.4, bt+0.05, 5.5, colw=[2.7,1.3,1.5], fs=9.5, hfs=9.5, rh=0.40)
notebox(s, "⚠ Caveat bắt buộc", "Spec/Prec 100% tựa trên CHỈ 2 bệnh nhân không-u → rất yếu thống kê; là phân biệt lát nội-bệnh-nhân, KHÔNG phải specificity lâm sàng.",
        0.4, bt+2.05, 5.5, 0.95, kind="caveat")
image(s, "seminar/code/detect_before_after.png", 6.1, bt+0.2, w=3.6,
      caption="Trước → sau fine-tune")
footer(s, 11)
notes(s, "Câu nói: Bắt đúng 20/25 bệnh nhân, FP thấp. Nhưng 100% kia đứng trên 2 người — chúng em không dùng nó để claim specificity lâm sàng.\nZero-shot gần như luôn trả 'No liver tumor is found' → fine-tune dạy được task từ số 0.")

# S10
s = newslide(); bt = header(s, "Kết quả: Định vị + vì sao model SÓT u", "Phần C — Kết quả")
bullets(s, [
  ("**Localization TRUNG BÌNH** (đã tách): per-patient **pIoU ≈ 0.27**, CI95 [0.19, 0.36]; recall@0.25 ≈ **54%** (nhỏ 35% / vừa 54% / lớn 74%).", 0),
  ("**Sót KHÔNG chỉ do u nhỏ:** theo TỪNG ổ, ổ sót nhỏ hơn ~2× VÀ tương phản kém hơn; nhưng theo BỆNH NHÂN thì cỡ gần bằng nhau.", 0),
  ("**Low-contrast/isointense** (u sáng gần bằng gan) mới phân biệt được cả u to — pid116 u 5.63% vẫn sót.", 1),
  ("Nghi cửa sổ CT rộng (WL40/WW400) → **cửa sổ gan hẹp** là hướng thử tiếp.", 1),
], 0.4, bt, 5.5, 3.4, base=11)
image(s, "seminar/code/pid116_contrast.png", 6.15, bt+0.05, w=3.5, caption="pid116: u 5.63% vẫn khó thấy (isointense)")
image(s, "seminar/code/missed_patients_v3.png", 6.15, bt+2.35, w=3.5, caption="tổng quan các ca bị sót")
footer(s, 12)
notes(s, "Câu nói: Định vị mới trung bình, chúng em không giấu. Ca sót nghiêng về u nhỏ VÀ tương phản thấp — cả hai cùng góp.\n\nisointense = đồng tỉ trọng. pIoU = IoU có phạt khi dư/thiếu box. 'kích thước' = diện tích box 2D (proxy), không phải thể tích thật.")

# S11
s = newslide(); bt = header(s, "Kết quả: Fine-tune KHÔNG làm “quên”", "Phần C — Kết quả")
bullets(s, [
  ("Lo ngại: fine-tune hẹp có làm mất khả năng nói chuyện? → kiểm bằng **hội thoại đa lượt**, lưu JSON.", 0),
  ("Model vừa **khoanh u** vừa: trả JSON đúng key, làm thơ, tính 17×23=**391**, giải thích song ngữ.", 0),
  ("**Vì sao giữ được:** LoRA giữ nguyên trọng số gốc → thêm kỹ năng mà không xoá kỹ năng cũ.", 0),
  ("⚠️ Kiểm **ĐỊNH TÍNH** (vài ca), chưa phải benchmark tổng quát có số.", 1),
], 0.4, bt, 4.6, 3.4, base=11.5)
image(s, "seminar/code/chat_test_multiturn.png", 5.15, bt+0.35, w=4.55,
      caption="hội thoại đa lượt (dữ liệu thật)")
footer(s, 13)
notes(s, "Câu nói: Model vừa khoanh được u, vừa theo lệnh chính xác, giải thích song ngữ, làm thơ, tính 17×23=391. Nhờ LoRA giữ nguyên trọng số gốc — thêm kỹ năng mà không quên. Kiểm định tính, chưa đo MMLU nên chưa dám nói bảo toàn hoàn toàn.")

# =========================================================
# PHẦN D — SO SÁNH & TƯƠNG LAI
# =========================================================
# S12
s = newslide(); bt = header(s, "So sánh với paper gốc", "Phần D — KHÔNG so số trực tiếp; so tính chất")
table(s, ["Khía cạnh", "MedRegA gốc", "Của nhóm"], [
  ["Quy mô", "~40B, 16×H800", "~8B, 1×A100"],
  ["Vào 3D", "1 lát trung tâm", "đa lát + ca âm"],
  ["Nhãn vùng", "box gộp thô", "multi-box"],
  ["Đơn vị đánh giá", "per-image", "per-patient + CI"],
  ["Detect vs Localize", "gộp", "tách riêng"],
  ["Ca âm / FP · Bất định", "không", "FP + selective prediction"],
  ["Grounded report", "có", "chưa (→Future)"],
], 0.4, bt, 5.6, colw=[1.8,1.9,1.9], fs=8.8, hfs=9, rh=0.33)
notebox(s, "Định vị — đọc đúng cách", "Gốc IoU ~0.23 vs mình pIoU ~0.27 (cùng thang 0–1) — khác data/metric, KHÔNG so trực tiếp, KHÔNG kết luận hơn/kém. pIoU có PHẠT nên khắt khe hơn.",
        6.15, bt+0.1, 3.5, 1.5, kind="info")
notebox(s, "Đóng góp thật", "Không phải điểm số — mà là khung đánh giá trung thực + độ tin cậy, bổ sung cho benchmark của paper.",
        6.15, bt+1.75, 3.5, 1.1, kind="info")
footer(s, 14)
notes(s, "Câu nói: Chúng em không claim đánh bại MedRegA. Đóng góp là biến benchmark định vị thành workflow đánh giá gần lâm sàng: có/không u, khoanh ở đâu, khi nào không nên tin.")

# S13a
s = newslide(); bt = header(s, "Hướng tương lai (1): ảnh giàu hơn + đáng tin hơn", "Phần D — [ĐO ĐƯỢC] nối số Slide 9–10 · [GIẢ ĐỊNH] suy từ thiết kế")
table(s, ["#", "Hạn chế (nguồn)", "Hướng đi tiếp", "Chi phí"], [
  ["①", "Sót u isointense [ĐO ĐƯỢC]", "Multi-window (cửa sổ gan hẹp)", "🟢 rất thấp"],
  ["②", "3D→2D mất thông tin [GIẢ ĐỊNH]", "2.5D stack (z−1/z/z+1 vào RGB)", "🟢 thấp"],
  ["③", "LiTS đơn pha [GIẢ ĐỊNH]", "CT đa pha — trần của DỮ LIỆU", "🔴 cần data"],
  ["⑤", "2 ca âm, 1 dataset [ĐO ĐƯỢC]", "External validation + FROC + calibration", "🟡 vừa"],
], 0.4, bt+0.15, 9.2, colw=[0.4,3.4,3.9,1.5], fs=10, hfs=10, rh=0.52)
footer(s, 15)
notes(s, "Câu nói: Mạch 1 — làm ảnh giàu hơn (multi-window, 2.5D, đa pha) và kết quả đáng tin hơn (external validation, thêm ca âm). ①⑤ nối thẳng số Kết quả; ②③ suy từ thiết kế nên ghi rõ CHƯA đo.\n\nExternal validation = kiểm trên bộ CT khác (IRCADb/MSD). FROC = đường cong phát hiện đa tổn thương. calibration = hiệu chỉnh độ tự tin.")

# S13b
s = newslide(); bt = header(s, "Hướng tương lai (2): nói-có-dẫn-chứng · generalist", "Phần D — [ĐO ĐƯỢC] · [PHẠM VI] · [GIẢ ĐỊNH]")
table(s, ["#", "Hạn chế (nguồn)", "Hướng đi tiếp", "Chi phí"], [
  ["④", "pIoU chỉ 0.27 [ĐO ĐƯỢC]", "MedSAM khoanh khít + polygon-as-token", "🟡 vừa"],
  ["⑥", "Chỉ detect [PHẠM VI]", "Grounded report (mỗi câu neo 1 box)", "🟡 vừa"],
  ["⑦", "Chỉ u gan [GIẢ ĐỊNH]", "Generalist đa cơ quan", "🔴 dài hạn"],
  ["⑧", "Loss CE chỉ khớp token toạ độ [PHẠM VI]", "Cân bằng loss thông minh (uncertainty weighting + loss hình học GIoU)", "🟢 thấp"],
], 0.4, bt+0.15, 9.2, colw=[0.4,3.2,4.1,1.5], fs=10, hfs=10, rh=0.52)
footer(s)
notes(s, "Câu nói: Mạch 2 — từ chỉ-chỗ tiến tới nói-có-dẫn-chứng rồi generalist đa cơ quan; và loss thông minh hơn (hình học GIoU) thay vì chỉ khớp token toạ độ. Tên tool (MedSAM/FROC/IRCADb) cần tự xác minh nguồn.\n\nMedSAM = mô hình phân vùng ảnh y tế. polygon-as-token = viết đường viền u thành chuỗi toạ độ. GIoU = IoU tổng quát.")

# S14
s = newslide(); bt = header(s, "Kết luận", "Phần D")
bullets(s, [
  ("**Chứng minh một cơ chế:** region-centric (khoanh-rồi-đọc) **khả thi ở ~8B, 1 GPU**.", 0),
  ("**Đóng góp cốt lõi:** Detection F1 **0.89** · pIoU **0.27** · FP **1.5%** — nhưng giá trị nằm ở **khung đánh giá trung thực + selective prediction**, không phải điểm số.", 0),
  ("**Trung thực:** PoC trên 1 dataset/1 task; n nhỏ, CI rộng; không tuyên bố hệ thống lâm sàng.", 0),
  ("**Bản đồ mở rộng:** ảnh giàu hơn → đáng tin hơn → nói-có-dẫn-chứng & generalist.", 0),
], 0.4, bt, 9.2, 3.2, base=12.5)
notebox(s, "Câu chốt", "Chúng em không hứa đã tới đích — chỉ ra con đường và nói rõ chỗ nào còn dốc.",
        0.4, bt+3.05, 9.2, 0.6, kind="info")
footer(s, 17)
notes(s, "Câu chốt: Chúng em không claim đánh bại MedRegA về quy mô hay SOTA. Đóng góp chính là biến một benchmark định vị vùng thành workflow đánh giá gần lâm sàng hơn, và trung thực về giới hạn.")

out = "seminar/cai_tien/seminar_full_2026-07-05.pptx"
prs.save(out)
print("SAVED", out, "| n_slides =", len(prs.slides._sldIdLst))
